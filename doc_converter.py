# app.py — Document Converter Pro (Part 1/3)
# Część 1/3: Importy, konfiguracja, offline guard, diagnostyka, helpery, OCR, Ollama utils, session_state
#cos nowego
import streamlit as st
import io
import base64
import re
import requests
import logging
import os
import json
import socket
import ipaddress
import shutil
import subprocess
import platform
import importlib.util
import tempfile
from urllib.parse import urlparse
from datetime import datetime
from PIL import Image
import numpy as np
import tempfile

# Parsery/formaty
import pdfplumber
from pdf2image import convert_from_bytes
import pytesseract
from pptx import Presentation
from docx import Document
import cv2

from typing import List, Dict, Any

# Opcjonalne biblioteki (ładowane dynamicznie; jeśli brak, będą None — UI pokaże w diagnostyce)
try:
    import mailparser
except Exception:
    mailparser = None
try:
    import extract_msg
except Exception:
    extract_msg = None
try:
    from duckduckgo_search import DDGS
except Exception:
    DDGS = None
try:
    import trafilatura
except Exception:
    trafilatura = None
try:
    import pandas as pd
except Exception:
    pd = None
    
def init_dc_state():
    """Inicjalizuje stan sesji – wyniki będą trwałe między rerunami."""
    ss = st.session_state
    ss.setdefault("converted", False)  # czy mamy już gotowe wyniki
    ss.setdefault("results", [])  # lista dict: {name, text, original_text, meta, pages}
    ss.setdefault("audio_items", [])  # lista: (name, text, meta)
    ss.setdefault("audio_summaries", [])  # opcjonalne podsumowania
    ss.setdefault("stats", {'processed': 0, 'errors': 0, 'pages': 0})
    ss.setdefault("file_sig", None)  # sygnatura zestawu plików (nazwa+rozmiar)
    ss.setdefault("speaker_maps", {})  # mapy imion per plik: {file_key: {SPEAKER_00:"Michał", ...}}

def files_signature(files) -> int:
    """Sygnatura zestawu plików – pomaga nie przeliczać ponownie po rerunie."""
    try:
        items = [(f.name, getattr(f, 'size', None) or len(f.getvalue())) for f in files]
        return hash(tuple(items))
    except Exception:
        return 0

init_dc_state()
# === SESSION STATE INIT ===

# Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("doc-converter")

st.set_page_config(page_title="📄 Document Converter", layout="wide", page_icon="📄")

# === CONFIG (domyślnie localhost; offline mode ON) ===
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434")
ANYTHINGLLM_URL = os.getenv("ANYTHINGLLM_URL", "http://anythingllm:3001")  # puste = wyłączone
ANYTHINGLLM_API_KEY = os.getenv("ANYTHINGLLM_API_KEY", "M2CCA9E-ARP41Y7-PZKT1G3-CK5SV2B")
WHISPER_URL = os.getenv("WHISPER_URL", "http://127.0.0.1:9000")
PYANNOTE_URL = os.getenv("PYANNOTE_URL", "http://127.0.0.1:8000")
OFFLINE_MODE = os.getenv("STRICT_OFFLINE", "1").lower() in ("1", "true", "yes")
ALLOW_WEB = False  # UI może to zmienić (web lookup)

# === STAŁE ===
MIN_TEXT_FOR_OCR_SKIP = 100
VISION_TRANSCRIBE_PROMPT = (
    "Przepisz DOKŁADNIE cały tekst z obrazu. Zachowaj pisownię, układ, symbole. "
    "Nie tłumacz, nie interpretuj - tylko przepisz. Jeśli coś nieczytelne - wpisz [NIECZYTELNE]."
)
VISION_DESCRIBE_PROMPT = (
    "Opisz ten obraz: co na nim widać? Wymień kluczowe elementy, teksty, wykresy lub diagramy, "
    "ogólny kontekst i ewentualny przekaz."
)
IMAGE_MODE_MAP = {
    "OCR": "ocr",
    "Vision: przepisz tekst": "vision_transcribe",
    "Vision: opisz obraz": "vision_describe",
    "OCR + Vision opis": "ocr_plus_vision_desc",
}

import re

def remap_speakers(text_with_speakers: str, speaker_map: dict) -> str:
    """
    Zamienia etykiety 'SPEAKER_00'... na podane imiona/role.
    Działa bezpiecznie, nie modyfikuje innych fragmentów.
    """
    if not text_with_speakers or not speaker_map:
        return text_with_speakers
    
    out = text_with_speakers
    # sortuj po długości klucza malejąco (żeby 'SPEAKER_1' nie nadpisał 'SPEAKER_10')
    for old in sorted(speaker_map.keys(), key=len, reverse=True):
        new = speaker_map.get(old, "").strip()
        if new and new != old:
            out = out.replace(old, new)
    return out

def speaker_mapper_form(file_name: str, original_text: str, current_text: str) -> str | None:
    """
    Rysuje formularz do nadania imion mówcom dla jednego pliku.
    Zwraca nowy tekst (po remap) albo None, jeśli nie kliknięto submit.
    """
    file_key = safe_filename(file_name)
    
    # wyciągnij mówców z ORYGINALNEGO tekstu (zachowuje SPEAKER_XX)
    unique = sorted(list(set(re.findall(r"SPEAKER_\d+", original_text))))
    if not unique:
        return None
    
    # mapa imion w sesji
    if file_key not in st.session_state["speaker_maps"]:
        st.session_state["speaker_maps"][file_key] = {}
    smap = st.session_state["speaker_maps"][file_key]
    
    with st.form(key=f"form_map_{file_key}"):
        cols = st.columns(len(unique))
        for i, spk in enumerate(unique):
            with cols[i]:
                smap[spk] = st.text_input(
                    f"Imię dla {spk}",
                    value=smap.get(spk, spk),
                    key=f"in_{file_key}_{spk}"
                )
        submitted = st.form_submit_button("Zastosuj imiona")
        if submitted:
            # zapisz mapę do sesji, zbuduj nowy tekst po remap
            st.session_state["speaker_maps"][file_key] = {k: v.strip() or k for k in smap.keys() for v in [smap[k]]}
            new_text = remap_speakers(original_text, st.session_state["speaker_maps"][file_key])
            st.success(f"Zastosowano imiona dla: {file_name}")
            return new_text
    return None
# === OFFLINE GUARD ===
def is_private_host(host: str) -> bool:
    try:
        infos = socket.getaddrinfo(host, None)
        for info in infos:
            sockaddr = info[4]
            ip = sockaddr[0] if isinstance(sockaddr, tuple) else sockaddr
            try:
                ip_obj = ipaddress.ip_address(ip)
                if ip_obj.is_private or ip_obj.is_loopback or ip_obj.is_link_local:
                    return True
            except ValueError:
                continue
    except Exception:
        return False
    return False

def assert_private_url(url: str):
    # Blokuje wyjście w internet w trybie OFFLINE_MODE (wyjątek: prywatne/localhost)
    if not OFFLINE_MODE:
        return
    try:
        p = urlparse(url)
        host = p.hostname or ""
        if not host or is_private_host(host):
            return
    except Exception:
        pass
    raise RuntimeError(f"Zablokowano żądanie poza sieć lokalną: {url}")

def http_get(url, **kwargs):
    assert_private_url(url)
    return requests.get(url, **kwargs)

def http_post(url, **kwargs):
    assert_private_url(url)
    return requests.post(url, **kwargs)

# === HELPERY PLIKÓW / FORMATÓW ===
def safe_filename(name: str) -> str:
    base = os.path.basename(name)
    base = re.sub(r'[^A-Za-z0-9.-]+', '', base)
    return base or "plik"

def create_run_dir(base_dir: str) -> str:
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    run_dir = os.path.join(base_dir, f"run_{ts}")
    os.makedirs(run_dir, exist_ok=True)
    return run_dir

def save_text(path: str, text: str):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write(text or "")

def format_timestamp(seconds: float) -> str:
    ms = int(round((seconds - int(seconds)) * 1000))
    s = int(seconds)
    h = s // 3600
    s = s % 3600
    m = s // 60
    s = s % 60
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"

def segments_to_srt(segments: list) -> str:
    lines = []
    for i, seg in enumerate(segments, 1):
        start = format_timestamp(seg.get("start", 0.0))
        end = format_timestamp(seg.get("end", 0.0))
        text_seg = (seg.get("text") or "").strip()
        lines.append(f"{i}\n{start} --> {end}\n{text_seg}\n")
    return "\n".join(lines)

# --- TIMEOUTY I HELPERY ROZMIARU ---
def get_file_size(file) -> int:
    try:
        sz = getattr(file, "size", None)
        if sz is not None:
            return sz
        try:
            return len(file.getvalue())
        except Exception:
            pass
        pos = file.tell()
        file.seek(0, os.SEEK_END)
        sz = file.tell()
        file.seek(pos)
        return sz
    except Exception:
        return 0

def calculate_timeout(file_size_bytes: int, base: int = 240, per_mb: int = 25) -> int:
    size_mb = max(1.0, file_size_bytes / 1024 / 1024)
    return int(max(base, size_mb * per_mb))

# === DIAGNOSTYKA ŚRODOWISKA ===
def has_module(mod: str) -> bool:
    return importlib.util.find_spec(mod) is not None

def module_version(mod: str) -> str:
    try:
        m = importlib.import_module(mod)
        return getattr(m, "__version__", "?")
    except Exception:
        return "?"

def cmd_version(cmd: str, args: list = ["--version"]) -> str:
    try:
        res = subprocess.run([cmd] + args, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, timeout=3)
        out = (res.stdout or "").strip().splitlines()
        return out[0] if out else "?"
    except Exception:
        return "not found"

def list_tesseract_langs() -> List[str]:
    try:
        res = subprocess.run(["tesseract", "--list-langs"], stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, timeout=4)
        out = res.stdout or ""
        langs = [l.strip() for l in out.splitlines() if l.strip() and not l.lower().startswith("list of available languages")]
        return langs
    except Exception:
        return []

def probe_service(name: str, base_url: str, paths: List[str]) -> Dict[str, Any]:
    info = {"name": name, "url": base_url, "ok": False, "detail": ""}
    try:
        for p in paths:
            url = base_url.rstrip("/") + p
            try:
                r = http_get(url, timeout=3)
                if r.ok:
                    info["ok"] = True
                    info["detail"] = f"OK {p} ({r.status_code})"
                    return info
                else:
                    info["detail"] = f"{p} -> HTTP {r.status_code}"
            except Exception as e:
                info["detail"] = f"{p} -> {e}"
        return info
    except Exception as e:
        info["detail"] = str(e)
        return info
def run_diagnostics() -> Dict[str, Any]:
    py_mods = [
        "pdfplumber", "pdf2image", "pytesseract", "pptx", "docx", "cv2", "PIL",
        "numpy", "requests", "pydub", "duckduckgo_search", "trafilatura",
        "mailparser", "extract_msg", "pandas"
    ]
    modules = {m: {"present": has_module(m), "version": module_version(m) if has_module(m) else ""} for m in py_mods}

    # executables: dodaj AMD ROCm narzędzia
    exes = {
        "tesseract": shutil.which("tesseract") is not None,
        "ffmpeg": shutil.which("ffmpeg") is not None,
        "pdftoppm": shutil.which("pdftoppm") is not None,
        "pdftocairo": shutil.which("pdftocairo") is not None,
        "nvidia-smi": shutil.which("nvidia-smi") is not None,
        "rocm-smi": shutil.which("rocm-smi") is not None,
        "rocminfo": shutil.which("rocminfo") is not None,
    }
    versions = {
        "tesseract": cmd_version("tesseract"),
        "ffmpeg": cmd_version("ffmpeg"),
        "pdftoppm": cmd_version("pdftoppm", ["-v"]),
        "pdftocairo": cmd_version("pdftocairo", ["-v"]),
        "nvidia-smi": cmd_version("nvidia-smi"),
        "rocm-smi": cmd_version("rocm-smi", ["--showproductname"]),
        "rocminfo": cmd_version("rocminfo", []),
    }

    langs = list_tesseract_langs()
    has_pol = any(l.lower() in ("pol", "polish") for l in langs)
    has_eng = any(l.lower() in ("eng", "english") for l in langs)

    services = {
        "ollama": probe_service("ollama", OLLAMA_URL, ["/api/tags", "/"]),
        "whisper": probe_service("whisper", WHISPER_URL, ["/health", "/", "/status"]),
        "pyannote": probe_service("pyannote", PYANNOTE_URL, ["/health", "/status", "/ping"])
    }

    missing_sys = []
    if not exes["tesseract"]:  missing_sys.append("tesseract")
    if not exes["ffmpeg"]:     missing_sys.append("ffmpeg")
    if not exes["pdftoppm"]:   missing_sys.append("poppler-tools (pdftoppm)")
    missing_langs = []
    if not has_pol: missing_langs.append("tesseract-langpack-pol (pol)")
    if not has_eng: missing_langs.append("tesseract-langpack-eng (eng)")

    missing_py = [m for m, v in modules.items() if not v["present"]]

    rec_zypper = []
    if "tesseract" in missing_sys: rec_zypper.append("tesseract")
    if "ffmpeg" in missing_sys: rec_zypper.append("ffmpeg")
    if "poppler-tools (pdftoppm)" in missing_sys: rec_zypper.append("poppler-tools")
    rec_zypper += [pkg for pkg in ["tesseract-langpack-pol", "tesseract-langpack-eng"] if pkg in [x.split()[0] for x in missing_langs]]

    rec_pip = missing_py if missing_py else []

    # GPU info: preferuj AMD (ROCm), jeśli rocm-smi jest dostępne
    gpu = {}
    if exes["rocm-smi"]:
        try:
            # Przykład: pokaż nazwę produktu i użycie pamięci
            prod = subprocess.run(["rocm-smi", "--showproductname"], stdout=subprocess.PIPE, text=True, timeout=4).stdout.strip()
            memu = subprocess.run(["rocm-smi", "--showmemuse"], stdout=subprocess.PIPE, text=True, timeout=4).stdout.strip()
            gpu["vendor"] = "AMD ROCm"
            gpu["product"] = prod
            gpu["mem_use"] = memu
        except Exception as e:
            gpu["vendor"] = "AMD ROCm"
            gpu["info"] = f"rocm-smi error: {e}"
    elif exes["nvidia-smi"]:
        try:
            res = subprocess.run(["nvidia-smi", "--query-gpu=name,memory.total", "--format=csv,noheader"], stdout=subprocess.PIPE, text=True, timeout=4)
            gpu["vendor"] = "NVIDIA"
            gpu["info"] = res.stdout.strip()
        except Exception as e:
            gpu["vendor"] = "NVIDIA"
            gpu["info"] = f"nvidia-smi error: {e}"
    else:
        gpu["vendor"] = "unknown"
        gpu["info"] = "GPU tool not found (no rocm-smi/nvidia-smi)"

    return {
        "system": {
            "os": platform.platform(),
            "executables": exes,
            "versions": versions,
            "tesseract_langs": langs,
            "gpu": gpu
        },
        "python_modules": modules,
        "services": services,
        "missing": {
            "system_packages": missing_sys,
            "tesseract_languages": missing_langs,
            "python_modules": missing_py
        },
        "recommend_install": {
            "zypper": rec_zypper,
            "pip": rec_pip
        }
    }

# === OLLAMA / MODELE / WIZJA ===
def list_ollama_models():
    try:
        r = http_get(f"{OLLAMA_URL}/api/tags", timeout=5)
        if r.ok:
            return [m.get("name", "") for m in r.json().get("models", [])]
    except Exception as e:
        logger.error(f"Ollama connection error: {e}")
    return []

def list_vision_models():
    all_models = list_ollama_models()
    prefixes = ("llava", "bakllava", "moondream", "llava-phi", "qwen2-vl")
    return [m for m in all_models if any(m.startswith(p) for p in prefixes)]

def query_ollama_vision(prompt: str, image_b64: str, model: str):
    try:
        payload = {"model": model, "prompt": prompt, "images": [image_b64], "stream": False}
        r = http_post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=120)
        r.raise_for_status()
        return r.json().get("response", "")
    except Exception as e:
        logger.error(f"Vision model error: {e}")
        return f"[BŁĄD VISION: {e}]"

def query_ollama_text(prompt: str, model: str = "llama3:latest", json_mode: bool = False, timeout: int = 120) -> str:
    try:
        payload = {"model": model, "prompt": prompt, "stream": False}
        if json_mode:
            payload["format"] = "json"
        r = http_post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=timeout)
        r.raise_for_status()
        return r.json().get("response", "")
    except Exception as e:
        logger.error(f"Ollama text error: {e}")
        return f"[BŁĄD OLLAMA: {e}]"

# === OCR (Tesseract z prostym preprocessingiem) ===
def ocr_image_bytes(img_bytes: bytes, lang: str = 'pol+eng') -> str:
    try:
        img = Image.open(io.BytesIO(img_bytes)).convert('L')
        np_img = np.array(img)
        _, thr = cv2.threshold(np_img, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        img_thr = Image.fromarray(thr)
        return pytesseract.image_to_string(img_thr, lang=lang) or ""
    except Exception as e:
        logger.warning(f"OCR error: {e}")
        return ""



# KONIEC CZĘŚCI 1/3 — daj znać, kiedy wysłać Część 2/3 (ekstraktory, audio/pyannote, podsumowania, Project Brain)
# app.py — Document Converter Pro (Part 2/3)
# Część 2/3: Ekstraktory plików, audio/pyannote, router, AnythingLLM,
#            podsumowania audio oraz Project Brain (zadania/ryzyka/brief + web lookup)

# === AUDIO: Whisper / Pyannote ===
def extract_audio_whisper(file):
    """Audio → tekst przez Whisper ASR. Zwraca (text, pages, meta)."""
    try:
        size_bytes = get_file_size(file)
        timeout_read = calculate_timeout(size_bytes, base=240, per_mb=25)

        file.seek(0)
        raw = file.read()
        mime = getattr(file, "type", None) or "application/octet-stream"
        fname = file.name

        # Opcjonalne downsample do 16k mono dla dużych plików (jeśli pydub/ffmpeg dostępne)
        try:
            if size_bytes > 25 * 1024 * 1024:
                from pydub import AudioSegment
                with tempfile.NamedTemporaryFile(suffix=os.path.splitext(fname)[1], delete=False) as tmp_in:
                    tmp_in.write(raw)
                    tmp_in_path = tmp_in.name
                audio = AudioSegment.from_file(tmp_in_path)
                audio = audio.set_channels(1).set_frame_rate(16000)
                buf = io.BytesIO()
                audio.export(buf, format="wav", parameters=["-acodec", "pcm_s16le"])
                raw = buf.getvalue()
                mime = "audio/wav"
                fname = os.path.splitext(fname)[0] + "_16k.wav"
                try:
                    os.remove(tmp_in_path)
                except Exception:
                    pass
        except Exception:
            pass

        files = {"audio_file": (fname, raw, mime)}
        r = http_post(
            f"{WHISPER_URL}/asr?task=transcribe&language=pl&word_timestamps=false&output=json",
            files=files,
            timeout=(30, timeout_read)
        )
        r.raise_for_status()

        try:
            result = r.json()
        except json.JSONDecodeError as je:
            logger.error(f"Whisper JSON decode error: {je}, response: {r.text[:500]}")
            return "[BŁĄD: Whisper zwrócił nieprawidłowy format]", 0, {"type": "audio", "error": "invalid_json"}

        text_res = result.get("text", "") or ""
        segments = result.get("segments", [])
        meta = {
            "type": "audio",
            "segments_count": len(segments),
            "duration": result.get("duration"),
            "language": result.get("language"),
            "segments": segments
        }

        if segments:
            lines = ["=== TRANSKRYPCJA Z TIMESTAMPAMI ===", ""]
            for seg in segments:
                start = seg.get("start", 0)
                end = seg.get("end", 0)
                txt = (seg.get("text") or "").strip()
                lines.append(f"[{start:.1f}s - {end:.1f}s] {txt}")
            text_res = "\n".join(lines)

        return text_res, 1, meta

    except requests.exceptions.Timeout:
        logger.error("Whisper timeout")
        return "[BŁĄD: Timeout - plik zbyt długi]", 0, {"type": "audio", "error": "timeout"}
    except Exception as e:
        logger.error(f"Whisper error: {e}")
        return f"[BŁĄD AUDIO: {e}]", 0, {"type": "audio", "error": str(e)}

def check_pyannote_health(url: str):
    url = (url or "").rstrip("/")
    for path in ("/health", "/status", "/ping"):
        try:
            r = http_get(url + path, timeout=3)
            if r.ok:
                try:
                    js = r.json()
                except Exception:
                    js = {"raw": r.text}
                if "model_loaded" in js:
                    return bool(js.get("model_loaded")), js
                return True, js
        except Exception:
            continue
    return False, {}

def normalize_diarization(resp: dict) -> list:
    """Ujednolicenie odpowiedzi z Pyannote do listy segmentów {start, end, speaker}."""
    if not resp:
        return []
    raw = resp.get("segments") or resp.get("turns") or []
    out = []
    for seg in raw:
        start = seg.get("start") or seg.get("start_time") or seg.get("begin") or 0.0
        end = seg.get("end") or seg.get("end_time") or seg.get("stop") or 0.0
        speaker = seg.get("speaker") or seg.get("label") or "SPEAKER_?"
        out.append({"start": float(start), "end": float(end), "speaker": speaker})
    return out

def pick_speaker_for_interval(diar_segments: list, start: float, end: float) -> str:
    """Wybierz mówcę z największym overlapem dla [start, end]."""
    best_spk, best_overlap = "SPEAKER_?", 0.0
    for s in diar_segments:
        ov = max(0.0, min(end, s["end"]) - max(start, s["start"]))
        if ov > best_overlap:
            best_overlap = ov
            best_spk = s["speaker"]
    return best_spk

def diarize_audio(file):
    """
    Wysyła plik audio do serwera pyannote, konwertując go w locie do WAV 16kHz mono.
    """
    pyannote_url = os.getenv("PYANNOTE_URL", PYANNOTE_URL).rstrip("/")
    
    # Zapisz oryginalny plik tymczasowo na dysku, aby ffmpeg mógł go odczytać
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_in:
        file.seek(0)
        tmp_in.write(file.read())
        original_audio_path = tmp_in.name

    converted_audio_path = None
    try:
        # --- POCZĄTEK TWOJEGO KODU (zintegrowany) ---
        
        # Jeśli już jest WAV, możemy pominąć konwersję (choć resamplowanie może być nadal potrzebne)
        # Dla pewności, zawsze konwertujemy do 16kHz mono
        with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_out:
            converted_audio_path = tmp_out.name
        
        logger.info(f"Konwertowanie pliku audio do WAV 16kHz mono dla pyannote...")
        command = [
            'ffmpeg', '-i', original_audio_path,
            '-ar', '16000',  # 16kHz sample rate
            '-ac', '1',       # mono channel
            '-y',             # overwrite output file if it exists
            converted_audio_path
        ]
        
        subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True
        )
        logger.info(f"Konwersja zakończona. Plik tymczasowy: {converted_audio_path}")
        
        # --- KONIEC TWOJEGO KODU ---

        # Wyślij skonwertowany plik do serwera pyannote
        with open(converted_audio_path, "rb") as converted_file:
            size_bytes = os.path.getsize(converted_audio_path)
            timeout_read = calculate_timeout(size_bytes, base=300, per_mb=30)
            
            files = {'file': (os.path.basename(converted_audio_path), converted_file, 'audio/wav')}
            
            r = http_post(
                f"{pyannote_url}/diarize",
                files=files,
                timeout=(30, timeout_read)
            )
            r.raise_for_status()
            return r.json()

    except subprocess.CalledProcessError as e:
        logger.error(f"Błąd ffmpeg podczas konwersji audio: {e.stderr}")
        return {"error": f"Błąd ffmpeg: {e.stderr}"}
    except requests.exceptions.ReadTimeout:
        logger.error("Timeout serwera pyannote.")
        return {"error": "timeout"}
    except Exception as e:
        logger.error(f"Błąd podczas diaryzacji: {e}")
        return {"error": str(e)}
    finally:
        # Posprzątaj pliki tymczasowe
        if os.path.exists(original_audio_path):
            os.remove(original_audio_path)
        if converted_audio_path and os.path.exists(converted_audio_path):
            os.remove(converted_audio_path)

def extract_audio_with_speakers(file):
    """Whisper + Pyannote = transkrypcja z identyfikacją głosów."""
    try:
        # 1) Whisper
        text_only, _, meta = extract_audio_whisper(file)
        segments = meta.get("segments", [])
        if not segments:
            return text_only, 1, meta

        # 2) Health-check Pyannote
        ok, _ = check_pyannote_health(PYANNOTE_URL)
        if not ok:
            st.warning("Nie udało się połączyć z Pyannote — zwracam samą transkrypcję.")
            return text_only, 1, meta

        # 3) Diarization
        st.info("🎤 Identyfikuję głosy...")
        file.seek(0)
        diarization = diarize_audio(file)
        diar_segments = normalize_diarization(diarization)
        if not diar_segments:
            st.warning("Pyannote nie zwrócił poprawnych segmentów — zwracam samą transkrypcję.")
            return text_only, 1, meta

        # 4) Merge
        output_lines = ["=== TRANSKRYPCJA Z IDENTYFIKACJĄ GŁOSÓW ===", ""]
        for seg in segments:
            start = float(seg.get("start", 0))
            end = float(seg.get("end", 0))
            txt = (seg.get("text") or "").strip()
            speaker = pick_speaker_for_interval(diar_segments, start, end)
            output_lines.append(f"[{start:.1f}s - {end:.1f}s] {speaker}: {txt}")

        result = "\n".join(output_lines)
        meta['has_speakers'] = True
        return result, 1, meta

    except Exception as e:
        logger.error(f"Audio with speakers error: {e}")
        return f"[BŁĄD: {e}]", 0, {"type": "audio", "error": str(e)}

# === EKSTRAKTORY DOKUMENTÓW ===
def extract_pdf(file, use_vision: bool, vision_model: str, ocr_pages_limit: int = 20):
    """PDF: tekst + opcjonalnie OCR/Vision (transkrypcja obrazu)."""
    texts = []
    try:
        file.seek(0)
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages):
                if i >= ocr_pages_limit:
                    texts.append(f"\n[... limit {ocr_pages_limit} stron ...]")
                    break
                t = page.extract_text() or ""
                texts.append(t)

        full_text = "\n".join(texts)

        # Fallback: OCR/Vision dla skanów lub niskiej jakości PDF
        if len(full_text.strip()) < MIN_TEXT_FOR_OCR_SKIP:
            file.seek(0)
            raw = file.read()
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                total_pages = len(pdf.pages)

            images = convert_from_bytes(
                raw,
                fmt="jpeg",
                dpi=150,
                first_page=1,
                last_page=min(ocr_pages_limit, total_pages)
            )

            if use_vision and vision_model:
                st.info(f"🖼️ Używam {vision_model} do analizy obrazów...")
                for idx, img in enumerate(images[:ocr_pages_limit], 1):
                    st.caption(f"Przetwarzam stronę {idx}/{min(ocr_pages_limit, len(images))}")
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG", quality=85)
                    img_b64 = base64.b64encode(buf.getvalue()).decode()
                    response = query_ollama_vision(VISION_TRANSCRIBE_PROMPT, img_b64, vision_model)
                    texts.append(f"\n--- Strona {idx} (Vision) ---\n{response}")
            else:
                st.info("📝 OCR Tesseract...")
                for idx, img in enumerate(images[:ocr_pages_limit], 1):
                    st.caption(f"OCR strona {idx}/{min(ocr_pages_limit, len(images))}")
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG")
                    ocr_text = ocr_image_bytes(buf.getvalue())
                    texts.append(f"\n--- Strona {idx} (OCR) ---\n{ocr_text}")

        meta = {"type": "pdf", "pages": len(texts)}
        return "\n".join(texts), len(texts), meta
    except Exception as e:
        logger.error(f"PDF extract error: {e}")
        return f"[BŁĄD PDF: {e}]", 0, {"type": "pdf", "error": str(e)}

def extract_pptx(file, use_vision: bool, vision_model: str):
    """PPTX: tekst + obrazy (opcjonalnie Vision opis)."""
    try:
        file.seek(0)
        prs = Presentation(file)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            parts = [f"=== Slajd {i} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    parts.append(f"Notatki: {notes}")

            if use_vision and vision_model:
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) == 13:  # PICTURE
                        try:
                            img_stream = shape.image.blob
                            img_b64 = base64.b64encode(img_stream).decode()
                            response = query_ollama_vision(VISION_DESCRIBE_PROMPT, img_b64, vision_model)
                            parts.append(f"[Obraz] {response}")
                        except Exception:
                            pass

            slides_text.append("\n".join(parts))

        return "\n\n".join(slides_text), len(prs.slides), {"type": "pptx", "slides": len(prs.slides)}
    except Exception as e:
        logger.error(f"PPTX error: {e}")
        return f"[BŁĄD PPTX: {e}]", 0, {"type": "pptx", "error": str(e)}

def extract_docx(file):
    """DOCX: tekst + tabele."""
    try:
        file.seek(0)
        doc = Document(file)
        paras = [p.text for p in doc.paragraphs if p.text]

        for tbl in doc.tables:
            for row in tbl.rows:
                paras.append(" | ".join(cell.text for cell in row.cells))

        return "\n".join(paras), len(paras), {"type": "docx"}
    except Exception as e:
        logger.error(f"DOCX error: {e}")
        return f"[BŁĄD DOCX: {e}]", 0, {"type": "docx", "error": str(e)}

def extract_image(file, use_vision: bool, vision_model: str, image_mode: str):
    """Obraz: OCR / Vision (przepisz) / Vision (opisz) / OCR+opis."""
    try:
        file.seek(0)
        img_bytes = file.read()
        results = []
        meta = {"type": "image", "mode": image_mode}

        if image_mode in ("ocr", "ocr_plus_vision_desc"):
            ocr_text = ocr_image_bytes(img_bytes)
            results.append(f"=== OCR ===\n{ocr_text}")

        if image_mode in ("vision_transcribe", "vision_describe", "ocr_plus_vision_desc"):
            if use_vision and vision_model:
                img_b64 = base64.b64encode(img_bytes).decode()
                prompt = VISION_TRANSCRIBE_PROMPT if image_mode == "vision_transcribe" else VISION_DESCRIBE_PROMPT
                vis = query_ollama_vision(prompt, img_b64, vision_model)
                tag = "Vision (transkrypcja)" if image_mode == "vision_transcribe" else "Vision (opis)"
                results.append(f"=== {tag} ===\n{vis}")
                meta["vision_model"] = vision_model
            else:
                results.append("[Vision niedostępne]")

        txt = "\n\n".join(results).strip()
        return txt, 1, meta
    except Exception as e:
        logger.error(f"Image error: {e}")
        return f"[BŁĄD IMG: {e}]", 0, {"type": "image", "error": str(e)}

# === E-MAIL: EML/MSG ===
def extract_eml(file):
    """EML: nagłówki + treść. Wymaga mailparser (opcjonalnie fallback)."""
    try:
        file.seek(0)
        raw = file.read()
        if mailparser is None:
            text = raw.decode("utf-8", errors="ignore")
            return text, 1, {"type": "email", "note": "mailparser not installed"}
        m = mailparser.parse_from_bytes(raw)
        headers = [
            f"From: {m.from_[0][1] if m.from_ else ''}",
            f"To: {', '.join([x[1] for x in m.to]) if m.to else ''}",
            f"Subject: {m.subject or ''}",
            f"Date: {m.date.isoformat() if m.date else ''}"
        ]
        body = (m.text_plain[0] if m.text_plain else m.body) or ""
        attach = [att.get('filename') for att in (m.attachments or []) if att.get('filename')]
        if attach:
            headers.append(f"Attachments: {', '.join(attach)}")
        text = "\n".join(headers) + "\n\n" + (body or "")
        return text, 1, {"type": "email", "has_attachments": bool(attach)}
    except Exception as e:
        logger.error(f"EML error: {e}")
        return f"[BŁĄD EML: {e}]", 0, {"type": "email", "error": str(e)}

def parse_msg_email(file):
    """MSG (Outlook): nagłówki + treść. Wymaga extract_msg (fallback na surowy tekst)."""
    try:
        if extract_msg is None:
            file.seek(0)
            raw = file.read()
            text = raw.decode("utf-8", errors="ignore")
            return text, 1, {"type": "email", "note": "extract-msg not installed"}
        file.seek(0)
        raw = file.read()
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        msg = extract_msg.Message(tmp_path)
        headers = [
            f"From: {msg.sender or ''}",
            f"To: {msg.to or ''}",
            f"Subject: {msg.subject or ''}",
            f"Date: {msg.date or ''}"
        ]
        body = msg.body or ""
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        text = "\n".join(headers) + "\n\n" + body
        return text, 1, {"type": "email"}
    except Exception as e:
        logger.error(f"MSG error: {e}")
        return f"[BŁĄD MSG: {e}]", 0, {"type": "email", "error": str(e)}

# === ROUTER ===
def process_file(file, use_vision: bool, vision_model: str, ocr_limit: int, image_mode: str):
    """Router do odpowiedniego ekstraktora."""
    name = file.name.lower()

    if name.endswith('.pdf'):
        return extract_pdf(file, use_vision, vision_model, ocr_limit)
    elif name.endswith(('.pptx', '.ppt')):
        return extract_pptx(file, use_vision, vision_model)
    elif name.endswith('.docx'):
        return extract_docx(file)
    elif name.endswith(('.jpg', '.jpeg', '.png')):
        return extract_image(file, use_vision, vision_model, image_mode)
    elif name.endswith(('.mp3', '.wav', '.m4a', '.ogg', '.flac')):
        ok, _ = check_pyannote_health(PYANNOTE_URL)
        if ok:
            return extract_audio_with_speakers(file)
        return extract_audio_whisper(file)
    elif name.endswith('.eml'):
        return extract_eml(file)
    elif name.endswith('.msg'):
        return parse_msg_email(file)
    elif name.endswith('.txt'):
        file.seek(0)
        content = file.read().decode('utf-8', errors='ignore')
        return content, 1, {"type": "txt"}
    else:
        return "[Nieobsługiwany format]", 0, {"type": "unknown"}

# === ANYTHINGLLM (lokalnie, respektuje offline) ===
def send_to_anythingllm(text: str, filename: str):
    """Wyślij dokument do AnythingLLM (tylko gdy nie OFFLINE_MODE i URL prywatny)."""
    if OFFLINE_MODE:
        return False, "Tryb offline — wysyłka zablokowana"
    if not ANYTHINGLLM_URL or not ANYTHINGLLM_API_KEY:
        return False, "Brak konfiguracji AnythingLLM"
    try:
        headers = {"Authorization": f"Bearer {ANYTHINGLLM_API_KEY}"}
        payload = {"name": filename, "content": text, "type": "text/plain"}
        r = http_post(f"{ANYTHINGLLM_URL}/api/v1/document-upload", headers=headers, json=payload, timeout=30)
        r.raise_for_status()
        return True, "✅ Wysłano do AnythingLLM"
    except Exception as e:
        return False, f"Błąd AnythingLLM: {e}"

# === PODSUMOWANIA AUDIO (Map-Reduce JSON + Markdown) ===
MAP_PROMPT_TEMPLATE = """
Jesteś asystentem ds. spotkań (PL). Otrzymasz fragment transkrypcji rozmowy z klientem (możliwe znaczniki SPEAKER_1, SPEAKER_2 i znaczniki czasu).
Zrób skrót tego fragmentu i wylistuj najważniejsze informacje.

WYMAGANY JSON:
{{
  "summary": "1-2 akapity skrótu (PL)",
  "key_points": ["punkt 1", "punkt 2", "..."],
  "decisions": ["decyzja 1", "decyzja 2"],
  "to_be_decided": ["kwestia do ustalenia 1", "kwestia 2"],
  "action_items": [{{ "owner":"", "task":"", "due":"", "notes":"" }}],
  "risks": [{{ "risk":"", "impact":"niski/średni/wysoki", "mitigation":"" }}],
  "open_questions": ["pytanie 1", "pytanie 2"]
}}

ZASADY:
- Nie wymyślaj informacji. Jeśli czegoś brak, zostaw puste pola lub wpisz [].
- Zostaw język polski.
- "to_be_decided" zawiera kwestie wymagające decyzji lub doprecyzowania.
- Jeśli są mówcy (SPEAKER_x) – staraj się zmapować właścicieli zadań (owner).
- Nie dodawaj komentarzy poza JSON.
Fragment:
{fragment}
"""

REDUCE_PROMPT_TEMPLATE = """
Jesteś asystentem ds. spotkań (PL). Otrzymasz listę częściowych podsumowań w JSON (z pól: summary, key_points, decisions, to_be_decided, action_items, risks, open_questions).
Scal je i zwróć jeden końcowy JSON w tym samym formacie. Usuń duplikaty, uczyść i pogrupuj logicznie.

WYMAGANY JSON:
{{
  "summary": "skondensowany skrót całości",
  "key_points": [...],
  "decisions": [...],
  "to_be_decided": [...],
  "action_items": [...],
  "risks": [...],
  "open_questions": [...]
}}

Wejście (lista JSON fragmentów):
{partials}

Nie dodawaj komentarzy poza JSON.
"""

def chunk_text(text: str, max_chars: int = 6000, overlap: int = 500) -> List[str]:
    if not text:
        return []
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + max_chars, n)
        chunks.append(text[start:end])
        if end >= n:
            break
        start = max(0, end - overlap)
    return chunks

def try_parse_json(s: str) -> Dict[str, Any]:
    if not s:
        return {}
    clean = s.strip()
    if clean.startswith("```json"):
        clean = clean[7:]
    if clean.startswith("```"):
        clean = clean[3:]
    if clean.endswith("```"):
        clean = clean[:-3]
    try:
        return json.loads(clean)
    except Exception:
        return {}

def merge_summary_dicts(items: List[Dict[str, Any]]) -> Dict[str, Any]:
    out = {
        "summary": "",
        "key_points": [],
        "decisions": [],
        "to_be_decided": [],
        "action_items": [],
        "risks": [],
        "open_questions": []
    }
    for it in items:
        if not isinstance(it, dict):
            continue
        if it.get("summary"):
            out["summary"] += (("\n" if out["summary"] else "") + it["summary"])
        for k in ["key_points", "decisions", "to_be_decided", "open_questions"]:
            out[k].extend(it.get(k, []))
        for ai in it.get("action_items", []):
            if isinstance(ai, dict):
                out["action_items"].append(ai)
        for r in it.get("risks", []):
            if isinstance(r, dict):
                out["risks"].append(r)
    for k in ["key_points", "decisions", "to_be_decided", "open_questions"]:
        out[k] = list(dict.fromkeys(out[k]))
    return out

def build_meeting_summary_markdown(data: Dict[str, Any]) -> str:
    if not data:
        return "_Brak danych do podsumowania_"
    md = []
    md.append("# Podsumowanie rozmowy")
    if data.get("summary"):
        md.append(data["summary"])

    md.append("\n## Kluczowe punkty")
    for x in (data.get("key_points") or []) or ["brak"]:
        md.append(f"- {x}")

    md.append("\n## Decyzje vs Do ustalenia")
    md.append("### Decyzje")
    for d in (data.get("decisions") or []) or ["brak"]:
        md.append(f"- {d}")
    md.append("### Do ustalenia")
    for q in (data.get("to_be_decided") or []) or ["brak"]:
        md.append(f"- {q}")

    md.append("\n## Zadania (Action Items)")
    action_items = data.get("action_items", [])
    if not action_items:
        md.append("- brak")
    else:
        for ai in action_items:
            owner = ai.get("owner", "") or "N/A"
            task = ai.get("task", "") or "N/A"
            due = ai.get("due", "") or "-"
            notes = ai.get("notes", "") or ""
            md.append(f"- [ ] {task} (owner: {owner}, termin: {due}) {('- ' + notes) if notes else ''}")

    md.append("\n## Ryzyka")
    risks = data.get("risks", [])
    if not risks:
        md.append("- brak")
    else:
        for r in risks:
            risk = r.get("risk", "")
            impact = r.get("impact", "")
            mit = r.get("mitigation", "")
            md.append(f"- {risk} (wpływ: {impact}) → mitygacja: {mit}")

    md.append("\n## Pytania do klienta (otwarte kwestie)")
    for q in (data.get("open_questions") or []) or ["brak"]:
        md.append(f"- {q}")

    return "\n".join(md)

def summarize_meeting_transcript(transcript: str, model: str = "llama3:latest", max_chars: int = 6000, diarized: bool = False) -> Dict[str, Any]:
    if not transcript or len(transcript.strip()) < 20:
        return {}

    parts = chunk_text(transcript, max_chars=max_chars, overlap=500)
    partials: List[Dict[str, Any]] = []
    for p in parts:
        prompt = MAP_PROMPT_TEMPLATE.format(fragment=p)
        resp = query_ollama_text(prompt, model=model, json_mode=True, timeout=180)
        data = try_parse_json(resp)
        if not data:
            resp2 = query_ollama_text(prompt, model=model, json_mode=False, timeout=180)
            data = try_parse_json(resp2)
        if data:
            partials.append(data)

    if not partials:
        return {"summary": transcript[:1200] + ("..." if len(transcript) > 1200 else "")}

    partials_str = json.dumps(partials, ensure_ascii=False, indent=2)
    reduce_prompt = REDUCE_PROMPT_TEMPLATE.format(partials=partials_str)
    reduce_resp = query_ollama_text(reduce_prompt, model=model, json_mode=True, timeout=240)
    final_data = try_parse_json(reduce_resp)
    if not final_data:
        final_data = merge_summary_dicts(partials)
    return final_data

# === PROJECT BRAIN (klasyfikacja, zadania, ryzyka, brief) ===
DOC_CLASS_PROMPT = """
Oceń, jaki to typ dokumentu (PL). Zwróć JSON z polem "type" ∈
["email","chat","meeting_transcript","spec","invoice","drawing","image_text","note","other"].
Kontekst:
{content}
Zwróć: {{"type": "<...>"}}
"""

TASKS_PROMPT = """
Jesteś asystentem PM (PL). Z treści wyodrębnij listę zadań.
Zwróć JSON:
{
 "tasks":[
   {"owner":"","task":"","due":"","priority":"low/medium/high","tags":[],"source":""}
 ]
}
Zasady: nie wymyślaj; jeśli brak due/owner wpisz "".
Treść:
{content}
"""

RISKS_PROMPT = """
Jesteś asystentem PM (PL). Wylistuj ryzyka, założenia i pytania (RFI).
JSON:
{
 "risks":[{"risk":"","impact":"low/medium/high","mitigation":""}],
 "assumptions":["..."],
 "rfis":["pytanie 1","pytanie 2"]
}
Treść:
{content}
"""

PROJECT_BRIEF_PROMPT = """
Z wielu fragmentów projektu zrób skrót (PL).
Wejście (JSON):
{items}
Zwróć JSON:
{
 "brief":"1-3 akapity",
 "key_points":[],
 "decisions":[],
 "rfis":[],
 "risks":[{"risk":"","impact":"","mitigation":""}],
 "next_steps":[]
}
"""

WEB_QUERIES_PROMPT = """
Na bazie treści zaproponuj 3–5 neutralnych zapytań do wyszukiwarki, bez danych wrażliwych.
JSON: {"queries":["...","..."]}
Treść:
{content}
"""

def classify_document(text: str, model: str) -> str:
    resp = query_ollama_text(DOC_CLASS_PROMPT.format(content=text[:4000]), model=model, json_mode=True, timeout=90)
    data = try_parse_json(resp)
    return (data.get("type") or "other") if isinstance(data, dict) else "other"

def extract_tasks_from_text(text: str, model: str) -> List[Dict[str,Any]]:
    resp = query_ollama_text(TASKS_PROMPT.format(content=text[:6000]), model=model, json_mode=True, timeout=120)
    data = try_parse_json(resp) or {}
    tasks = data.get("tasks", []) if isinstance(data, dict) else []
    # sanityzacja typów
    out = []
    for t in tasks:
        if isinstance(t, dict):
            out.append({
                "owner": t.get("owner",""),
                "task": t.get("task",""),
                "due": t.get("due",""),
                "priority": t.get("priority",""),
                "tags": t.get("tags", []),
                "source": t.get("source","")
            })
    return out

def extract_risks_from_text(text: str, model: str) -> Dict[str,Any]:
    resp = query_ollama_text(RISKS_PROMPT.format(content=text[:6000]), model=model, json_mode=True, timeout=120)
    data = try_parse_json(resp) or {}
    return {
        "risks": data.get("risks", []),
        "assumptions": data.get("assumptions", []),
        "rfis": data.get("rfis", [])
    }

def build_project_brief(items: List[Dict[str,Any]], model: str) -> Dict[str,Any]:
    payload = json.dumps(items, ensure_ascii=False)[:12000]
    resp = query_ollama_text(PROJECT_BRIEF_PROMPT.format(items=payload), model=model, json_mode=True, timeout=180)
    data = try_parse_json(resp)
    if not data:
        return {"brief": "", "key_points": [], "decisions": [], "rfis": [], "risks": [], "next_steps": []}
    return data

def propose_web_queries(text: str, model: str) -> List[str]:
    resp = query_ollama_text(WEB_QUERIES_PROMPT.format(content=text[:4000]), model=model, json_mode=True, timeout=90)
    data = try_parse_json(resp) or {}
    q = data.get("queries", [])
    # proste odanonimizowanie (usuń maile/telefony)
    clean = []
    for s in q:
        s = re.sub(r'\S+@\S+', '[email]', s)
        s = re.sub(r'\b\d{7,}\b', '[num]', s)
        clean.append(s)
    return clean

def web_search_and_summarize(queries: List[str], max_results: int, model: str) -> Dict[str,Any]:
    """Pobiera strony przez DuckDuckGo i streszcza lokalnie.
       Działa tylko, gdy w UI włączono ALLOW_WEB (sesja)."""
    if not st.session_state.get("ALLOW_WEB", False):
        return {"note":"web lookup disabled (ALLOW_WEB=False)", "items":[]}
    if DDGS is None or trafilatura is None:
        return {"note":"duckduckgo-search/trafilatura not installed", "items":[]}
    results = []
    try:
        with DDGS() as ddg:
            for q in queries or []:
                hits = ddg.text(q, region="pl-pl", safesearch="moderate", max_results=max_results)
                for h in (hits or []):
                    url = h.get("href") or h.get("url")
                    title = h.get("title","")
                    if not url:
                        continue
                    # pobierz treść
                    try:
                        c = trafilatura.fetch_url(url)
                        txt = trafilatura.extract(c) or ""
                    except Exception:
                        txt = ""
                    if not txt:
                        continue
                    summary = query_ollama_text(f"Streść (PL) w 5 punktach:\n\n{txt[:6000]}", model=model, json_mode=False, timeout=90)
                    results.append({"query": q, "url": url, "title": title, "summary": summary})
                    if len(results) >= max_results:
                        break
                if len(results) >= max_results:
                    break
    except Exception as e:
        logger.error(f"Web lookup error: {e}")
    return {"items": results}

# KONIEC CZĘŚCI 2/3 — daj znać, kiedy wysłać Część 3/3 (UI: sidebar, uploader, konwersja, zapisy, Project Brain UI, web lookup)
# app.py — Document Converter Pro (Part 3/3)
# Część 3/3: UI — sidebar, uploader, konwersja, zapisy, podsumowania audio,
#            Project Brain (zadania/ryzyka/brief) + opcjonalny web lookup

# === UI / SIDEBAR ===
st.title("📄 Document Converter Pro")
st.caption("Konwersja PDF/DOCX/PPTX/IMG/AUDIO/EMAIL → TXT z OCR, Vision lub Whisper (offline-first)")

with st.sidebar:
    st.header("⚙️ Ustawienia")

    # Tryb offline
    OFFLINE_MODE = st.checkbox("Tryb offline (blokuj internet poza lokalnymi usługami)", value=OFFLINE_MODE)
    st.session_state["ALLOW_WEB"] = st.checkbox(
        "Zezwól na web lookup (pobieranie publicznych stron)",
        value=st.session_state.get("ALLOW_WEB", False),
        help="Nie wysyła treści dokumentów na zewnątrz. Pobiera tylko publiczne strony dla uzupełnienia wiedzy."
    )

    # Status adresów
    def _status_url(name, url):
        try:
            host = urlparse(url).hostname or ""
            st.caption(f"{name}: {url} → {'✅ lokalny/prywatny' if is_private_host(host) else '❌ zewnętrzny'}")
        except Exception:
            st.caption(f"{name}: {url} → ⚠️ nie można zweryfikować")
    _status_url("Ollama", OLLAMA_URL)
    _status_url("Whisper", WHISPER_URL)
    _status_url("Pyannote", PYANNOTE_URL)

    # Vision
    vision_models = list_vision_models()
    use_vision = st.checkbox("Użyj modelu wizyjnego (Ollama Vision)", value=True if vision_models else False)
    if vision_models and use_vision:
        selected_vision = st.selectbox("Model wizyjny", vision_models, index=0)
    else:
        selected_vision = None
        if use_vision:
            st.warning("⚠️ Brak modeli Vision w Ollama (np. llava:13b / qwen2-vl:7b)")

    st.subheader("OCR")
    ocr_pages_limit = st.slider("Limit stron OCR", 5, 50, 20)

    st.subheader("Obrazy (IMG)")
    if use_vision and selected_vision:
        image_mode_label = st.selectbox(
            "Tryb dla obrazów",
            options=["OCR", "Vision: przepisz tekst", "Vision: opisz obraz", "OCR + Vision opis"],
            index=3
        )
    else:
        image_mode_label = st.selectbox(
            "Tryb dla obrazów",
            options=["OCR"],
            index=0,
            disabled=True
        )
    IMAGE_MODE_MAP = {
        "OCR": "ocr",
        "Vision: przepisz tekst": "vision_transcribe",
        "Vision: opisz obraz": "vision_describe",
        "OCR + Vision opis": "ocr_plus_vision_desc",
    }
    image_mode = IMAGE_MODE_MAP.get(image_mode_label, "ocr")

    # Zapis lokalny
    st.subheader("Zapis lokalny")
    enable_local_save = st.checkbox("Zapisz wyniki lokalnie (folder)", value=False)
    base_output_dir = st.text_input("Katalog wyjściowy", value="outputs")

    # AnythingLLM
    st.subheader("AnythingLLM")
    has_anythingllm_cfg = bool(ANYTHINGLLM_URL and ANYTHINGLLM_API_KEY)
    if OFFLINE_MODE:
        st.caption("Status: 🔒 Wyłączone (tryb offline)")
        has_anythingllm = False
    else:
        st.caption(f"Status: {'✅ Skonfigurowane' if has_anythingllm_cfg else '❌ Brak config'}")
        has_anythingllm = has_anythingllm_cfg

    # Podsumowania audio
    st.subheader("🧠 Podsumowanie audio (AI)")
    summarize_audio_enabled = st.checkbox("Włącz podsumowanie rozmów audio", value=True)
    summarize_model_candidates = [
        m for m in list_ollama_models()
        if not any(m.startswith(p) for p in ("llava", "bakllava", "moondream", "llava-phi", "nomic-embed", "qwen2-vl"))
    ]
    summarize_model = st.selectbox("Model do podsumowania", options=summarize_model_candidates or ["llama3:latest"])
    chunk_chars = st.slider("Rozmiar chunku (znaki)", min_value=2000, max_value=8000, value=6000, step=500)

    # Project Brain toggle
    st.subheader("🧭 Project Brain (PM)")
    enable_project_brain = st.checkbox("Włącz Project Brain (zadania/ryzyka/brief)", value=True)

    # Diagnostyka
    st.subheader("🔧 Diagnostyka środowiska")
    if st.button("Skanuj środowisko"):
        st.session_state["diag"] = run_diagnostics()
    if st.session_state.get("diag"):
        diag = st.session_state["diag"]
        st.caption("Wyniki diagnostyki (skrót):")
        miss = diag.get("missing", {})
        st.write(f"- Brakujące systemowe: {', '.join(miss.get('system_packages', []) ) or '—'}")
        st.write(f"- Brakujące języki Tesseract: {', '.join(miss.get('tesseract_languages', []) ) or '—'}")
        st.write(f"- Brakujące moduły Python: {', '.join(miss.get('python_modules', []) ) or '—'}")
        rec = diag.get("recommend_install", {})
        if rec.get("zypper") or rec.get("pip"):
            st.markdown("Zalecane instalacje:")
            if rec.get("zypper"):
                st.code("sudo zypper in -y " + " ".join(sorted(set(rec["zypper"])) ), language="bash")
            if rec.get("pip"):
                st.code("pip install " + " ".join(sorted(set(rec["pip"])) ), language="bash")
        with st.expander("Pełne szczegóły diagnostyki"):
            st.json(diag, expanded=False)

# === FILE UPLOADER ===
uploaded_files = st.file_uploader(
    "Wgraj dokumenty",
    type=['pdf','docx','pptx','ppt','jpg','jpeg','png','txt','mp3','wav','m4a','ogg','flac'],
    accept_multiple_files=True
)
# === KONWERSJA → zapis do session_state (bez resetu przy zapisie) ===
if uploaded_files:
    st.info(f"📁 {len(uploaded_files)} plików")
    if st.button("🚀 Konwertuj wszystkie", type="primary", key="btn_convert_all"):
        # Reset stanu dla nowego przebiegu
        st.session_state["results"] = []
        st.session_state["combined_text"] = ""
        st.session_state["audio_items"] = []
        st.session_state["audio_summaries"] = []
        st.session_state["stats"] = {'processed': 0, 'errors': 0, 'pages': 0}
        st.session_state["converted"] = False
        st.session_state["files_sig"] = files_signature(uploaded_files)
        st.session_state["run_dir"] = create_run_dir(base_output_dir) if enable_local_save else None
        if enable_local_save:
            st.info(f"💾 Wyniki będą zapisane w: {st.session_state['run_dir']}")

        progress = st.progress(0)
        all_texts = []

        for idx, file in enumerate(uploaded_files):
            try:
                progress.progress((idx + 1) / len(uploaded_files), text=f"Przetwarzam: {file.name}")
            except TypeError:
                progress.progress((idx + 1) / len(uploaded_files))

            st.subheader(f"📄 {file.name}")

            try:
                extracted_text, pages, meta = process_file(file, use_vision, selected_vision, ocr_pages_limit, image_mode)

                # Do sesji
                st.session_state["results"].append({
                    "name": file.name,
                    "text": extracted_text,
                    "original_text": extracted_text, # <-- DODAJ TĘ LINIĘ
                    "meta": meta,
                    "pages": pages
                })
               

                # Tekst łączny
                all_texts.append(f"\n{'='*80}\n")
                all_texts.append(f"PLIK: {file.name}\n")
                all_texts.append(f"Typ: {getattr(file, 'type', 'unknown')}, Rozmiar: {getattr(file, 'size', 0)/1024:.1f} KB\n")
                all_texts.append(f"{'='*80}\n")
                all_texts.append(extracted_text)
                all_texts.append(f"\n[Stron/sekcji: {pages}]\n")

                # Statystyki
                st.session_state["stats"]["processed"] += 1
                st.session_state["stats"]["pages"] += pages

                # Podgląd
                with st.expander(f"Preview: {file.name}"):
                    st.text(extracted_text[:2000] + ("..." if len(extracted_text) > 2000 else ""))


                # Audio → do podsumowań
                if isinstance(meta, dict) and meta.get("type") == "audio":
                    st.session_state["audio_items"].append((file.name, extracted_text, meta))

            except Exception as e:
                st.error(f"❌ Błąd: {e}")
                logger.exception(f"Error processing {file.name}")
                st.session_state["stats"]["errors"] += 1

        progress.empty()
        st.session_state["combined_text"] = "\n".join(all_texts)
        st.session_state["converted"] = True

# === SEKCJA WYNIKÓW ===
if st.session_state.get("converted"):
    st.success(f"✅ Przetworzono: {st.session_state['stats']['processed']} plików | Sekcji: {st.session_state['stats']['pages']}")
    
    # Wyświetl wyniki z możliwością edycji imion mówców
    for res in st.session_state["results"]:
        name = res["name"]
        text = res["text"]
        meta = res["meta"]
        
        st.subheader(f"📄 {name}")
        with st.expander("Podgląd", expanded=True):
            st.text_area(f"prev_{safe_filename(name)}", text, height=240, key=f"preview_{safe_filename(name)}")
        
        # FORMULARZ DO NADAWANIA IMION – nie wywołuje rerunu podczas pisania
        if isinstance(meta, dict) and meta.get("has_speakers"):
            st.markdown("##### 🎤 Przypisz imiona mówcom")
            new_text = speaker_mapper_form(name, res.get("original_text", text), text)
            if new_text is not None:
                # tylko gdy klikniesz Zastosuj – podmiana tekstu i rerun
                res["text"] = new_text
                st.rerun()
    
    # Przyciski eksportu
    st.download_button(
        "⬇️ Pobierz wszystko w TXT",
        ("\n".join([f"\n{'='*60}\nPLIK: {r['name']}\n{'='*60}\n{r['text']}" for r in st.session_state['results']])).encode("utf-8"),
        file_name=f"converted_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
        mime="text/plain",
        key="dl_all_txt"
    )

    c1, c2, c3, c4 = st.columns(4)

    with c1:
        if st.button("💾 Zapisz połączony TXT na dysk", key="btn_save_combined_txt"):
            out_dir = st.session_state["run_dir"] or create_run_dir("outputs")
            st.session_state["run_dir"] = out_dir
            combined_path = os.path.join(out_dir, f"converted_{datetime.now().strftime('%Y%m%d_%H%M')}.txt")
            save_text(combined_path, st.session_state["combined_text"])
            st.success(f"Zapisano: {combined_path}")

    with c2:
        if st.button("💾 Zapisz wszystkie SRT (audio)", key="btn_save_all_srt"):
            out_dir = st.session_state["run_dir"] or create_run_dir("outputs")
            st.session_state["run_dir"] = out_dir
            saved = 0
            for it in st.session_state["results"]:
                meta = it.get("meta") or {}
                if meta.get("type") == "audio" and meta.get("segments"):
                    fname_base = os.path.splitext(safe_filename(it["name"]))[0]
                    srt_path = os.path.join(out_dir, f"{fname_base}.srt")
                    save_text(srt_path, segments_to_srt(meta["segments"]))
                    saved += 1
            st.success(f"Zapisano SRT dla {saved} plików audio w: {out_dir}")

    with c3:
        if st.button("🧠 Generuj podsumowania audio (MD+JSON)", key="btn_make_summaries"):
            st.session_state["audio_summaries"] = []
            if st.session_state["audio_items"]:
                for (aname, atext, ameta) in st.session_state["audio_items"]:
                    diarized = bool(ameta.get("has_speakers"))
                    with st.spinner(f"Tworzę podsumowanie dla {aname}..."):
                        summary_json = summarize_meeting_transcript(
                            transcript=atext,
                            model=summarize_model if summarize_model_candidates else "llama3:latest",
                            max_chars=chunk_chars,
                            diarized=diarized
                        )
                        summary_md = build_meeting_summary_markdown(summary_json)
                        st.session_state["audio_summaries"].append({"name": aname, "md": summary_md, "json": summary_json})
                st.success("Gotowe podsumowania – poniżej do pobrania/zapisania.")
            else:
                st.info("Brak plików audio do podsumowania.")

    with c4:
        if has_anythingllm:
            if st.button("📤 Wyślij do AnythingLLM", key="btn_send_anythingllm"):
                success, msg = send_to_anythingllm(st.session_state["combined_text"], "converted_docs.txt")
                if success:
                    st.success(msg)
                else:
                    st.error(msg)
        else:
            st.caption("AnythingLLM wyłączone lub brak config/tryb offline")

    # === WYŚWIETLANIE PODSUMOWAŃ AUDIO ===
    if st.session_state.get("audio_summaries"):
        st.markdown("---")
        st.subheader("🧠 Podsumowania rozmów audio")
        
        for s in st.session_state["audio_summaries"]:
            aname = s["name"]
            summary_md = s["md"]
            summary_json = s["json"]

            st.markdown(f"### 🎧 {aname}")
            st.markdown(summary_md)

            col_a, col_b, col_c = st.columns(3)
            with col_a:
                st.download_button(
                    "⬇️ Pobierz MD",
                    summary_md.encode("utf-8"),
                    file_name=f"{safe_filename(aname).replace('.mp3', '').replace('.wav', '')}_summary.md",
                    mime="text/markdown",
                    key=f"dl_md_{safe_filename(aname)}"
                )
            with col_b:
                st.download_button(
                    "⬇️ Pobierz JSON",
                    json.dumps(summary_json, ensure_ascii=False, indent=2).encode("utf-8"),
                    file_name=f"{safe_filename(aname).replace('.mp3', '').replace('.wav', '')}_summary.json",
                    mime="application/json",
                    key=f"dl_json_{safe_filename(aname)}"
                )
            with col_c:
                if st.button("💾 Zapisz (MD+JSON) na dysk", key=f"btn_save_sum_{safe_filename(aname)}"):
                    out_dir = st.session_state.get("run_dir") or create_run_dir("outputs")
                    st.session_state["run_dir"] = out_dir
                    base = safe_filename(aname).replace('.mp3', '').replace('.wav', '')
                    
                    md_path = os.path.join(out_dir, f"{base}_summary.md")
                    json_path = os.path.join(out_dir, f"{base}_summary.json")
                    
                    save_text(md_path, summary_md)
                    save_text(json_path, json.dumps(summary_json, ensure_ascii=False, indent=2))
                    
                    st.success(f"Zapisano podsumowanie do: {out_dir}")

# Reset sesji (nie rusza plików na dysku)
# Reset sesji (nie rusza plików na dysku)
st.markdown("---")
if st.button("♻️ Reset sesji (wyczyść wyniki)", type="secondary", key="btn_reset_session"):
    for k in ["results", "combined_text", "audio_items", "audio_summaries", "run_dir", "project_brain", "project_tasks"]:
        st.session_state[k] = [] if isinstance(st.session_state.get(k), list) else None
    st.session_state["stats"] = {'processed': 0, 'errors': 0, 'pages': 0}
    st.session_state["converted"] = False
    st.info("Wyczyszczono wyniki z pamięci sesji.")
