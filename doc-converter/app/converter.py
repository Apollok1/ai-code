# app.py ‚Äî Document Converter Pro (scalony, z podsumowaniami audio + dokument√≥w)

# Czƒô≈õƒá 1/3: Importy, konfiguracja, offline guard, diagnostyka, helpery, OCR, Ollama utils, session_state

import streamlit as st
import io
import base64
import re
import requests
import logging
import os
import json
import socket
import ipaddress
import shutil
import subprocess
import platform
import importlib.util
import tempfile
from urllib.parse import urlparse
from datetime import datetime
from PIL import Image
import numpy as np

# Parsery/formaty
import pdfplumber
from pdf2image import convert_from_bytes
import pytesseract
from pptx import Presentation
from docx import Document
import cv2

from typing import List, Dict, Any

# Opcjonalne biblioteki (≈Çadowane dynamicznie; je≈õli brak, bƒôdƒÖ None ‚Äî UI poka≈ºe w diagnostyce)
try:
    import mailparser
except Exception:
    mailparser = None
try:
    import extract_msg
except Exception:
    extract_msg = None
try:
    from duckduckgo_search import DDGS
except Exception:
    DDGS = None
try:
    import trafilatura
except Exception:
    trafilatura = None
try:
    import pandas as pd
except Exception:
    pd = None


def init_dc_state():
    """Inicjalizuje stan sesji ‚Äì wyniki bƒôdƒÖ trwa≈Çe miƒôdzy rerunami."""
    ss = st.session_state
    ss.setdefault("converted", False)        # czy mamy ju≈º gotowe wyniki
    ss.setdefault("results", [])             # lista dict: {name, text, original_text, meta, pages}
    ss.setdefault("audio_items", [])         # lista: (name, text, meta)
    ss.setdefault("audio_summaries", [])     # podsumowania AUDIO
    ss.setdefault("doc_summaries", [])       # podsumowania DOKUMENT√ìW (tekst/obraz)
    ss.setdefault("stats", {'processed': 0, 'errors': 0, 'pages': 0})
    ss.setdefault("file_sig", None)          # sygnatura zestawu plik√≥w (nazwa+rozmiar)
    ss.setdefault("speaker_maps", {})        # mapy imion per plik: {file_key: {SPEAKER_00:"Micha≈Ç", ...}}


def files_signature(files) -> int:
    """Sygnatura zestawu plik√≥w ‚Äì pomaga nie przeliczaƒá ponownie po rerunie."""
    try:
        items = [(f.name, getattr(f, 'size', None) or len(f.getvalue())) for f in files]
        return hash(tuple(items))
    except Exception:
        return 0


init_dc_state()

# Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("doc-converter")

st.set_page_config(page_title="üìÑ Document Converter", layout="wide", page_icon="üìÑ")

# === CONFIG (domy≈õlnie localhost; offline mode ON) ===
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434")
ANYTHINGLLM_URL = os.getenv("ANYTHINGLLM_URL", "http://anythingllm:3001")  # puste = wy≈ÇƒÖczone
ANYTHINGLLM_API_KEY = os.getenv("ANYTHINGLLM_API_KEY", "")
WHISPER_URL = os.getenv("WHISPER_URL", "http://127.0.0.1:9000")
PYANNOTE_URL = os.getenv("PYANNOTE_URL", "http://127.0.0.1:8000")
OFFLINE_MODE = os.getenv("STRICT_OFFLINE", "1").lower() in ("1", "true", "yes")
ALLOW_WEB = False  # UI mo≈ºe to zmieniƒá (web lookup)

# === STA≈ÅE ===
MIN_TEXT_FOR_OCR_SKIP = 100
VISION_TRANSCRIBE_PROMPT = (
    "Przepisz DOK≈ÅADNIE ca≈Çy tekst z obrazu. Zachowaj pisowniƒô, uk≈Çad, symbole, formatowanie. "
    "ZASADY:\n"
    "- Nie t≈Çumacz, nie interpretuj - tylko przepisz\n"
    "- Zachowaj podzia≈Ç na akapity i linie\n"
    "- Je≈õli co≈õ nieczytelne - wpisz [NIECZYTELNE]\n"
    "- Je≈õli nie ma tekstu - napisz [BRAK TEKSTU]\n"
    "- Przepisuj cyfry, daty, nazwy dok≈Çadnie jak sƒÖ\n"
    "Pisz TYLKO po polsku (lub w oryginalnym jƒôzyku je≈õli tekst nie jest polski)."
)
VISION_DESCRIBE_PROMPT = (
    "Przeprowad≈∫ szczeg√≥≈ÇowƒÖ analizƒô technicznƒÖ tego obrazu po polsku.\n\n"
    "STRUKTURA ODPOWIEDZI:\n"
    "1. TYP OBIEKTU: Podaj nazwƒô technicznƒÖ i kategoriƒô (np. ≈Ço≈ºysko kulkowe, silnik, narzƒôdzie)\n"
    "2. MATERIA≈Å: Okre≈õl z czego wykonany (stal, aluminium, plastik, drewno, itp.)\n"
    "3. BUDOWA: Wymie≈Ñ wszystkie widoczne czƒô≈õci sk≈Çadowe i ich rozmieszczenie\n"
    "4. KSZTA≈ÅT I WYMIARY: Opisz geometriƒô, proporcje, charakterystyczne cechy\n"
    "5. FUNKCJA: Do czego s≈Çu≈ºy ten obiekt\n"
    "6. ZASTOSOWANIE: Gdzie jest typowo u≈ºywany\n\n"
    "ZASADY:\n"
    "- U≈ºywaj precyzyjnej terminologii in≈ºynieryjnej i technicznej\n"
    "- Opisuj TYLKO to co faktycznie widzisz na obrazie\n"
    "- Nie zgaduj, nie domy≈õlaj siƒô - je≈õli czego≈õ nie widaƒá, nie wymy≈õlaj\n"
    "- Je≈õli widzisz tekst/napisy - przepisz je dok≈Çadnie\n"
    "Pisz wy≈ÇƒÖcznie po polsku."
)
IMAGE_MODE_MAP = {
    "OCR": "ocr",
    "Vision: przepisz tekst": "vision_transcribe",
    "Vision: opisz obraz": "vision_describe",
    "OCR + Vision opis": "ocr_plus_vision_desc",
}


def remap_speakers(text_with_speakers: str, speaker_map: dict) -> str:
    """
    Zamienia etykiety 'SPEAKER_00'... na podane imiona/role.
    Dzia≈Ça bezpiecznie, nie modyfikuje innych fragment√≥w.
    """
    if not text_with_speakers or not speaker_map:
        return text_with_speakers

    out = text_with_speakers
    # sortuj po d≈Çugo≈õci klucza malejƒÖco (≈ºeby 'SPEAKER_1' nie nadpisa≈Ç 'SPEAKER_10')
    for old in sorted(speaker_map.keys(), key=len, reverse=True):
        new = speaker_map.get(old, "").strip()
        if new and new != old:
            out = out.replace(old, new)
    return out


def speaker_mapper_form(file_name: str, original_text: str, current_text: str) -> str | None:
    """
    Rysuje formularz do nadania imion m√≥wcom dla jednego pliku.
    Zwraca nowy tekst (po remap) albo None, je≈õli nie klikniƒôto submit.
    """
    file_key = safe_filename(file_name)

    # wyciƒÖgnij m√≥wc√≥w z ORYGINALNEGO tekstu (zachowuje SPEAKER_XX)
    unique = sorted(list(set(re.findall(r"SPEAKER_\d+", original_text))))
    if not unique:
        return None

    # mapa imion w sesji
    if file_key not in st.session_state["speaker_maps"]:
        st.session_state["speaker_maps"][file_key] = {}
    smap = st.session_state["speaker_maps"][file_key]

    with st.form(key=f"form_map_{file_key}"):
        cols = st.columns(len(unique))
        for i, spk in enumerate(unique):
            with cols[i]:
                smap[spk] = st.text_input(
                    f"Imiƒô dla {spk}",
                    value=smap.get(spk, spk),
                    key=f"in_{file_key}_{spk}"
                )
        submitted = st.form_submit_button("Zastosuj imiona")
        if submitted:
            # zapisz mapƒô do sesji, zbuduj nowy tekst po remap
            st.session_state["speaker_maps"][file_key] = {
                k: v.strip() or k for k, v in smap.items()
            }
            new_text = remap_speakers(original_text, st.session_state["speaker_maps"][file_key])
            st.success(f"Zastosowano imiona dla: {file_name}")
            return new_text
    return None


# === OFFLINE GUARD ===
def is_private_host(host: str) -> bool:
    # Whitelist known local Docker service names
    LOCAL_SERVICES = {"whisper", "pyannote", "ollama", "anythingllm", "localhost", "127.0.0.1", "::1"}
    if host.lower() in LOCAL_SERVICES:
        return True

    # Check if hostname ends with .local (Docker internal)
    if host.endswith(".local"):
        return True

    try:
        infos = socket.getaddrinfo(host, None)
        for info in infos:
            sockaddr = info[4]
            ip = sockaddr[0] if isinstance(sockaddr, tuple) else sockaddr
            try:
                ip_obj = ipaddress.ip_address(ip)
                if ip_obj.is_private or ip_obj.is_loopback or ip_obj.is_link_local:
                    return True
            except ValueError:
                continue
    except Exception:
        return False
    return False


def assert_private_url(url: str):
    # Blokuje wyj≈õcie w internet w trybie OFFLINE_MODE (wyjƒÖtek: prywatne/localhost)
    if not OFFLINE_MODE:
        return
    try:
        p = urlparse(url)
        host = p.hostname or ""
        if not host or is_private_host(host):
            return
    except Exception:
        pass
    raise RuntimeError(f"Zablokowano ≈ºƒÖdanie poza sieƒá lokalnƒÖ: {url}")


def http_get(url, **kwargs):
    assert_private_url(url)
    return requests.get(url, **kwargs)


def http_post(url, **kwargs):
    assert_private_url(url)
    return requests.post(url, **kwargs)


# === HELPERY PLIK√ìW / FORMAT√ìW ===
def safe_filename(name: str) -> str:
    base = os.path.basename(name)
    base = re.sub(r'[^A-Za-z0-9.-]+', '', base)
    return base or "plik"


def create_run_dir(base_dir: str) -> str:
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    run_dir = os.path.join(base_dir, f"run_{ts}")
    os.makedirs(run_dir, exist_ok=True)
    return run_dir


def save_text(path: str, text: str):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write(text or "")


def format_timestamp(seconds: float) -> str:
    ms = int(round((seconds - int(seconds)) * 1000))
    s = int(seconds)
    h = s // 3600
    s = s % 3600
    m = s // 60
    s = s % 60
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def segments_to_srt(segments: list) -> str:
    lines = []
    for i, seg in enumerate(segments, 1):
        start = format_timestamp(seg.get("start", 0.0))
        end = format_timestamp(seg.get("end", 0.0))
        text_seg = (seg.get("text") or "").strip()
        lines.append(f"{i}\n{start} --> {end}\n{text_seg}\n")
    return "\n".join(lines)


# --- TIMEOUTY I HELPERY ROZMIARU ---
def get_file_size(file) -> int:
    try:
        sz = getattr(file, "size", None)
        if sz is not None:
            return sz
        try:
            return len(file.getvalue())
        except Exception:
            pass
        pos = file.tell()
        file.seek(0, os.SEEK_END)
        sz = file.tell()
        file.seek(pos)
        return sz
    except Exception:
        return 0


def calculate_timeout(file_size_bytes: int, base: int = 240, per_mb: int = 25) -> int:
    size_mb = max(1.0, file_size_bytes / 1024 / 1024)
    return int(max(base, size_mb * per_mb))


# === DIAGNOSTYKA ≈öRODOWISKA ===
def has_module(mod: str) -> bool:
    return importlib.util.find_spec(mod) is not None


def module_version(mod: str) -> str:
    try:
        m = importlib.import_module(mod)
        return getattr(m, "__version__", "?")
    except Exception:
        return "?"


def cmd_version(cmd: str, args: list = ["--version"]) -> str:
    try:
        res = subprocess.run([cmd] + args, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, timeout=3)
        out = (res.stdout or "").strip().splitlines()
        return out[0] if out else "?"
    except Exception:
        return "not found"


def list_tesseract_langs() -> List[str]:
    try:
        res = subprocess.run(["tesseract", "--list-langs"], stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, timeout=4)
        out = res.stdout or ""
        langs = [l.strip() for l in out.splitlines() if l.strip() and not l.lower().startswith("list of available languages")]
        return langs
    except Exception:
        return []


def probe_service(name: str, base_url: str, paths: List[str]) -> Dict[str, Any]:
    info = {"name": name, "url": base_url, "ok": False, "detail": ""}
    try:
        for p in paths:
            url = base_url.rstrip("/") + p
            try:
                r = http_get(url, timeout=3)
                if r.ok:
                    try:
                        js = r.json()
                    except Exception:
                        js = {"raw": r.text}
                    if "model_loaded" in js:
                        info["ok"] = bool(js.get("model_loaded"))
                        info["detail"] = f"OK {p} ({r.status_code})"
                        return info
                    info["ok"] = True
                    info["detail"] = f"OK {p} ({r.status_code})"
                    return info
                else:
                    info["detail"] = f"{p} -> HTTP {r.status_code}"
            except Exception as e:
                info["detail"] = f"{p} -> {e}"
        return info
    except Exception as e:
        info["detail"] = str(e)
        return info


def run_diagnostics() -> Dict[str, Any]:
    py_mods = [
        "pdfplumber", "pdf2image", "pytesseract", "pptx", "docx", "cv2", "PIL",
        "numpy", "requests", "pydub", "duckduckgo_search", "trafilatura",
        "mailparser", "extract_msg", "pandas"
    ]
    modules = {m: {"present": has_module(m), "version": module_version(m) if has_module(m) else ""} for m in py_mods}

    # executables: dodaj AMD ROCm narzƒôdzia
    exes = {
        "tesseract": shutil.which("tesseract") is not None,
        "ffmpeg": shutil.which("ffmpeg") is not None,
        "pdftoppm": shutil.which("pdftoppm") is not None,
        "pdftocairo": shutil.which("pdftocairo") is not None,
        "nvidia-smi": shutil.which("nvidia-smi") is not None,
        "rocm-smi": shutil.which("rocm-smi") is not None,
        "rocminfo": shutil.which("rocminfo") is not None,
    }
    versions = {
        "tesseract": cmd_version("tesseract"),
        "ffmpeg": cmd_version("ffmpeg"),
        "pdftoppm": cmd_version("pdftoppm", ["-v"]),
        "pdftocairo": cmd_version("pdftocairo", ["-v"]),
        "nvidia-smi": cmd_version("nvidia-smi"),
        "rocm-smi": cmd_version("rocm-smi", ["--showproductname"]),
        "rocminfo": cmd_version("rocminfo", []),
    }

    langs = list_tesseract_langs()
    has_pol = any(l.lower() in ("pol", "polish") for l in langs)
    has_eng = any(l.lower() in ("eng", "english") for l in langs)

    services = {
        "ollama": probe_service("ollama", OLLAMA_URL, ["/api/tags", "/"]),
        "whisper": probe_service("whisper", WHISPER_URL, ["/health", "/", "/status"]),
        "pyannote": probe_service("pyannote", PYANNOTE_URL, ["/health", "/status", "/ping"])
    }

    missing_sys = []
    if not exes["tesseract"]:
        missing_sys.append("tesseract")
    if not exes["ffmpeg"]:
        missing_sys.append("ffmpeg")
    if not exes["pdftoppm"]:
        missing_sys.append("poppler-tools (pdftoppm)")
    missing_langs = []
    if not has_pol:
        missing_langs.append("tesseract-langpack-pol (pol)")
    if not has_eng:
        missing_langs.append("tesseract-langpack-eng (eng)")

    missing_py = [m for m, v in modules.items() if not v["present"]]

    rec_zypper = []
    if "tesseract" in missing_sys:
        rec_zypper.append("tesseract")
    if "ffmpeg" in missing_sys:
        rec_zypper.append("ffmpeg")
    if "poppler-tools (pdftoppm)" in missing_sys:
        rec_zypper.append("poppler-tools")
    rec_zypper += [pkg for pkg in ["tesseract-langpack-pol", "tesseract-langpack-eng"] if pkg in [x.split()[0] for x in missing_langs]]

    rec_pip = missing_py if missing_py else []

    # GPU info: preferuj AMD (ROCm), je≈õli rocm-smi jest dostƒôpne
    gpu = {}
    if exes["rocm-smi"]:
        try:
            prod = subprocess.run(["rocm-smi", "--showproductname"], stdout=subprocess.PIPE, text=True, timeout=4).stdout.strip()
            memu = subprocess.run(["rocm-smi", "--showmemuse"], stdout=subprocess.PIPE, text=True, timeout=4).stdout.strip()
            gpu["vendor"] = "AMD ROCm"
            gpu["product"] = prod
            gpu["mem_use"] = memu
        except Exception as e:
            gpu["vendor"] = "AMD ROCm"
            gpu["info"] = f"rocm-smi error: {e}"
    elif exes["nvidia-smi"]:
        try:
            res = subprocess.run(
                ["nvidia-smi", "--query-gpu=name,memory.total", "--format=csv,noheader"],
                stdout=subprocess.PIPE,
                text=True,
                timeout=4,
            )
            gpu["vendor"] = "NVIDIA"
            gpu["info"] = res.stdout.strip()
        except Exception as e:
            gpu["vendor"] = "NVIDIA"
            gpu["info"] = f"nvidia-smi error: {e}"
    else:
        gpu["vendor"] = "unknown"
        gpu["info"] = "GPU tool not found (no rocm-smi/nvidia-smi)"

    return {
        "system": {
            "os": platform.platform(),
            "executables": exes,
            "versions": versions,
            "tesseract_langs": langs,
            "gpu": gpu
        },
        "python_modules": modules,
        "services": services,
        "missing": {
            "system_packages": missing_sys,
            "tesseract_languages": missing_langs,
            "python_modules": missing_py
        },
        "recommend_install": {
            "zypper": rec_zypper,
            "pip": rec_pip
        }
    }


# === OLLAMA / MODELE / WIZJA ===
def list_ollama_models():
    try:
        r = http_get(f"{OLLAMA_URL}/api/tags", timeout=5)
        if r.ok:
            return [m.get("name", "") for m in r.json().get("models", [])]
    except Exception as e:
        logger.error(f"Ollama connection error: {e}")
    return []


def list_vision_models():
    all_models = list_ollama_models()
    vision_prefixes = (
        "llava", "bakllava", "moondream", "llava-phi",
        "qwen2-vl", "qwen2.5vl", "qwen",  # warianty Qwen Vision
        "cogvlm", "internvl", "minicpm-v"
    )
    return [m for m in all_models if any(m.startswith(p) for p in vision_prefixes)]


def query_ollama_vision(prompt: str, image_b64: str, model: str):
    try:
        payload = {"model": model, "prompt": prompt, "images": [image_b64], "stream": False}
        r = http_post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=120)
        r.raise_for_status()
        return r.json().get("response", "")
    except Exception as e:
        logger.error(f"Vision model error: {e}")
        return f"[B≈ÅƒÑD VISION: {e}]"


def query_ollama_text(prompt: str, model: str = None, json_mode: bool = False, timeout: int = 120) -> str:
    # U≈ºyj wybranego modelu z session_state je≈õli nie podano
    if model is None:
        model = st.session_state.get("selected_main_text_model", "qwen2.5:7b")

    try:
        payload = {"model": model, "prompt": prompt, "stream": False}
        if json_mode:
            payload["format"] = "json"
        r = http_post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=timeout)
        r.raise_for_status()
        return r.json().get("response", "")
    except Exception as e:
        logger.error(f"Ollama text error: {e}")
        return f"[B≈ÅƒÑD OLLAMA: {e}]"


# === OCR (Tesseract z prostym preprocessingiem) ===
def ocr_image_bytes(img_bytes: bytes, lang: str = 'pol+eng') -> str:
    try:
        img = Image.open(io.BytesIO(img_bytes)).convert('L')
        np_img = np.array(img)
        _, thr = cv2.threshold(np_img, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        img_thr = Image.fromarray(thr)
        return pytesseract.image_to_string(img_thr, lang=lang) or ""
    except Exception as e:
        logger.warning(f"OCR error: {e}")
        return ""


# Czƒô≈õƒá 2/3: Ekstraktory plik√≥w, audio/pyannote, router, AnythingLLM,
#            podsumowania audio oraz Project Brain (zadania/ryzyka/brief + web lookup)

# === AUDIO: Whisper / Pyannote ===

def split_audio_into_chunks(input_path: str, chunk_minutes: int = 10):
    """
    Dzieli plik audio na mniejsze chunki o okre≈õlonej d≈Çugo≈õci.

    Args:
        input_path: ≈öcie≈ºka do pliku audio
        chunk_minutes: D≈Çugo≈õƒá ka≈ºdego chunka w minutach

    Returns:
        Lista tupli: [(chunk_path, start_offset_seconds), ...]
    """
    try:
        # Pobierz d≈Çugo≈õƒá pliku audio
        probe_cmd = [
            'ffprobe', '-v', 'error',
            '-show_entries', 'format=duration',
            '-of', 'default=noprint_wrappers=1:nokey=1',
            input_path
        ]
        result = subprocess.run(probe_cmd, capture_output=True, text=True, check=True)
        total_duration = float(result.stdout.strip())

        chunk_seconds = chunk_minutes * 60
        chunks = []

        # Je≈õli plik kr√≥tszy ni≈º chunk, zwr√≥ƒá go bez dzielenia
        if total_duration <= chunk_seconds:
            return [(input_path, 0.0)]

        # Dziel na chunki
        start = 0.0
        chunk_idx = 0
        while start < total_duration:
            chunk_path = tempfile.NamedTemporaryFile(suffix='.wav', delete=False).name
            duration = min(chunk_seconds, total_duration - start)

            # ffmpeg: wytnij fragment
            split_cmd = [
                'ffmpeg', '-i', input_path,
                '-ss', str(start),
                '-t', str(duration),
                '-ar', '16000',  # 16kHz
                '-ac', '1',      # mono
                '-y',
                chunk_path
            ]

            subprocess.run(split_cmd, capture_output=True, check=True)
            chunks.append((chunk_path, start))

            logger.info(f"Utworzono chunk {chunk_idx}: {start:.1f}s - {start+duration:.1f}s")

            start += duration
            chunk_idx += 1

        return chunks

    except Exception as e:
        logger.error(f"B≈ÇƒÖd podczas dzielenia audio: {e}")
        return [(input_path, 0.0)]  # Fallback: zwr√≥ƒá oryginalny plik


def extract_audio_whisper(file, enable_chunking=False, chunk_minutes=10):
    """
    Audio ‚Üí tekst przez Whisper ASR. Zwraca (text, pages, meta).

    Args:
        file: Plik audio
        enable_chunking: Czy w≈ÇƒÖczyƒá dzielenie na czƒô≈õci (dla d≈Çugich plik√≥w)
        chunk_minutes: D≈Çugo≈õƒá ka≈ºdego chunka w minutach (domy≈õlnie 10)
    """
    try:
        size_bytes = get_file_size(file)
        file.seek(0)
        raw = file.read()
        fname = file.name

        # Zapisz do pliku tymczasowego
        with tempfile.NamedTemporaryFile(suffix=os.path.splitext(fname)[1], delete=False) as tmp_in:
            tmp_in.write(raw)
            input_path = tmp_in.name

        # Konwersja do WAV 16kHz mono (dla sp√≥jno≈õci)
        converted_path = tempfile.NamedTemporaryFile(suffix='.wav', delete=False).name
        try:
            convert_cmd = [
                'ffmpeg', '-i', input_path,
                '-ar', '16000',
                '-ac', '1',
                '-y',
                converted_path
            ]
            subprocess.run(convert_cmd, capture_output=True, check=True)
        except Exception as e:
            logger.warning(f"Konwersja audio nie powiod≈Ça siƒô: {e}, u≈ºywam oryginalnego pliku")
            converted_path = input_path

        try:
            # Podzia≈Ç na chunki je≈õli w≈ÇƒÖczone
            if enable_chunking:
                chunks = split_audio_into_chunks(converted_path, chunk_minutes)
                logger.info(f"Podzielono audio na {len(chunks)} czƒô≈õci")
            else:
                chunks = [(converted_path, 0.0)]

            # Przetwarzanie ka≈ºdego chunka
            all_segments = []
            all_texts = []
            total_duration = 0.0

            for chunk_idx, (chunk_path, time_offset) in enumerate(chunks):
                logger.info(f"Przetwarzanie chunka {chunk_idx + 1}/{len(chunks)} (offset: {time_offset:.1f}s)")

                # Odczytaj chunk
                with open(chunk_path, 'rb') as cf:
                    chunk_data = cf.read()

                chunk_size = len(chunk_data)
                timeout_read = calculate_timeout(chunk_size, base=300, per_mb=30)

                # Wy≈õlij do Whisper
                files = {"audio_file": (f"chunk_{chunk_idx}.wav", chunk_data, "audio/wav")}
                r = http_post(
                    f"{WHISPER_URL}/asr?task=transcribe&language=pl&word_timestamps=false&output=json",
                    files=files,
                    timeout=(30, timeout_read)
                )
                r.raise_for_status()

                try:
                    result = r.json()
                except json.JSONDecodeError as je:
                    logger.error(f"Whisper JSON decode error: {je}")
                    continue

                # Przetw√≥rz segmenty z chunka
                segments = result.get("segments", [])
                chunk_duration = result.get("duration", 0.0)

                for seg in segments:
                    # Dodaj offset czasowy do timestamp√≥w
                    seg["start"] = seg.get("start", 0) + time_offset
                    seg["end"] = seg.get("end", 0) + time_offset
                    all_segments.append(seg)

                total_duration = max(total_duration, time_offset + chunk_duration)

                # Usu≈Ñ plik chunka je≈õli to nie oryginalny
                if chunk_path != converted_path:
                    try:
                        os.remove(chunk_path)
                    except Exception:
                        pass

            # Z≈Ç√≥≈º wyniki
            meta = {
                "type": "audio",
                "segments_count": len(all_segments),
                "duration": total_duration,
                "language": "pl",
                "segments": all_segments,
                "chunked": enable_chunking,
                "chunks_count": len(chunks) if enable_chunking else 1
            }

            if all_segments:
                lines = ["=== TRANSKRYPCJA Z TIMESTAMPAMI ===", ""]
                for seg in all_segments:
                    start = seg.get("start", 0)
                    end = seg.get("end", 0)
                    txt = (seg.get("text") or "").strip()
                    lines.append(f"[{start:.1f}s - {end:.1f}s] {txt}")
                text_res = "\n".join(lines)
            else:
                text_res = "[Brak transkrypcji]"

            return text_res, 1, meta

        finally:
            # Cleanup
            try:
                os.remove(input_path)
                if converted_path != input_path:
                    os.remove(converted_path)
            except Exception:
                pass

    except requests.exceptions.Timeout:
        logger.error("Whisper timeout")
        return "[B≈ÅƒÑD: Timeout - plik zbyt d≈Çugi]", 0, {"type": "audio", "error": "timeout"}
    except Exception as e:
        logger.error(f"Whisper error: {e}")
        return f"[B≈ÅƒÑD AUDIO: {e}]", 0, {"type": "audio", "error": str(e)}


def check_pyannote_health(url: str):
    url = (url or "").rstrip("/")
    for path in ("/health", "/status", "/ping"):
        try:
            r = http_get(url + path, timeout=3)
            if r.ok:
                try:
                    js = r.json()
                except Exception:
                    js = {"raw": r.text}
                if "model_loaded" in js:
                    return bool(js.get("model_loaded")), js
                return True, js
        except Exception:
            continue
    return False, {}


def normalize_diarization(resp: dict) -> list:
    """Ujednolicenie odpowiedzi z Pyannote do listy segment√≥w {start, end, speaker}."""
    if not resp:
        return []
    raw = resp.get("segments") or resp.get("turns") or []
    out = []
    for seg in raw:
        start = seg.get("start") or seg.get("start_time") or seg.get("begin") or 0.0
        end = seg.get("end") or seg.get("end_time") or seg.get("stop") or 0.0
        speaker = seg.get("speaker") or seg.get("label") or "SPEAKER_?"
        out.append({"start": float(start), "end": float(end), "speaker": speaker})
    return out


def pick_speaker_for_interval(diar_segments: list, start: float, end: float) -> str:
    """Wybierz m√≥wcƒô z najwiƒôkszym overlapem dla [start, end]."""
    best_spk, best_overlap = "SPEAKER_?", 0.0
    for s in diar_segments:
        ov = max(0.0, min(end, s["end"]) - max(start, s["start"]))
        if ov > best_overlap:
            best_overlap = ov
            best_spk = s["speaker"]
    return best_spk


def diarize_audio(file, enable_chunking=False, chunk_minutes=10):
    """
    Wysy≈Ça plik audio do serwera pyannote, konwertujƒÖc go w locie do WAV 16kHz mono.

    Args:
        file: Plik audio
        enable_chunking: Czy w≈ÇƒÖczyƒá dzielenie na czƒô≈õci (dla d≈Çugich plik√≥w)
        chunk_minutes: D≈Çugo≈õƒá ka≈ºdego chunka w minutach (domy≈õlnie 10)
    """
    pyannote_url = os.getenv("PYANNOTE_URL", PYANNOTE_URL).rstrip("/")

    # Zapisz plik do temp
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_in:
        file.seek(0)
        tmp_in.write(file.read())
        original_audio_path = tmp_in.name

    converted_audio_path = None
    try:
        # Konwersja do WAV 16kHz mono
        converted_audio_path = tempfile.NamedTemporaryFile(suffix='.wav', delete=False).name

        logger.info(f"Konwertowanie pliku audio do WAV 16kHz mono dla pyannote...")
        command = [
            'ffmpeg', '-i', original_audio_path,
            '-ar', '16000',  # 16kHz
            '-ac', '1',      # mono
            '-y',
            converted_audio_path
        ]

        subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True
        )
        logger.info(f"Konwersja zako≈Ñczona. Plik tymczasowy: {converted_audio_path}")

        # Podzia≈Ç na chunki je≈õli w≈ÇƒÖczone
        if enable_chunking:
            chunks = split_audio_into_chunks(converted_audio_path, chunk_minutes)
            logger.info(f"Podzielono audio na {len(chunks)} czƒô≈õci dla diaryzacji")
        else:
            chunks = [(converted_audio_path, 0.0)]

        # Przetwarzanie ka≈ºdego chunka
        all_segments = []

        for chunk_idx, (chunk_path, time_offset) in enumerate(chunks):
            logger.info(f"Diaryzacja chunka {chunk_idx + 1}/{len(chunks)} (offset: {time_offset:.1f}s)")

            with open(chunk_path, "rb") as chunk_file:
                size_bytes = os.path.getsize(chunk_path)
                timeout_read = calculate_timeout(size_bytes, base=300, per_mb=30)

                files = {'audio_file': (f"chunk_{chunk_idx}.wav", chunk_file, 'audio/wav')}

                r = http_post(
                    f"{pyannote_url}/diarize",
                    files=files,
                    timeout=(30, timeout_read)
                )
                r.raise_for_status()
                chunk_result = r.json()

            # Normalizuj i dodaj offset czasowy
            chunk_segments = normalize_diarization(chunk_result)
            for seg in chunk_segments:
                seg["start"] += time_offset
                seg["end"] += time_offset
                all_segments.append(seg)

            # Usu≈Ñ plik chunka je≈õli to nie oryginalny
            if chunk_path != converted_audio_path:
                try:
                    os.remove(chunk_path)
                except Exception:
                    pass

        # Zwr√≥ƒá w formacie zgodnym z normalize_diarization
        return {"segments": all_segments, "chunked": enable_chunking}

    except subprocess.CalledProcessError as e:
        logger.error(f"B≈ÇƒÖd ffmpeg podczas konwersji audio: {e.stderr}")
        return {"error": f"B≈ÇƒÖd ffmpeg: {e.stderr}"}
    except requests.exceptions.ReadTimeout:
        logger.error("Timeout serwera pyannote.")
        return {"error": "timeout"}
    except Exception as e:
        logger.error(f"B≈ÇƒÖd podczas diaryzacji: {e}")
        return {"error": str(e)}
    finally:
        if os.path.exists(original_audio_path):
            os.remove(original_audio_path)
        if converted_audio_path and os.path.exists(converted_audio_path):
            os.remove(converted_audio_path)


def extract_audio_with_speakers(file, enable_chunking=False, chunk_minutes=10):
    """
    Whisper + Pyannote = transkrypcja z identyfikacjƒÖ g≈Ços√≥w.

    Args:
        file: Plik audio
        enable_chunking: Czy w≈ÇƒÖczyƒá dzielenie na czƒô≈õci
        chunk_minutes: D≈Çugo≈õƒá chunka w minutach
    """
    try:
        text_only, _, meta = extract_audio_whisper(file, enable_chunking, chunk_minutes)
        segments = meta.get("segments", [])
        if not segments:
            return text_only, 1, meta

        ok, _ = check_pyannote_health(PYANNOTE_URL)
        if not ok:
            st.warning("Nie uda≈Ço siƒô po≈ÇƒÖczyƒá z Pyannote ‚Äî zwracam samƒÖ transkrypcjƒô.")
            return text_only, 1, meta

        st.info("üé§ Identyfikujƒô g≈Çosy...")
        file.seek(0)
        diarization = diarize_audio(file, enable_chunking, chunk_minutes)

        # diarize_audio teraz zwraca dict z kluczem "segments"
        if "error" in diarization:
            st.warning(f"B≈ÇƒÖd diaryzacji: {diarization['error']} ‚Äî zwracam samƒÖ transkrypcjƒô.")
            return text_only, 1, meta

        diar_segments = diarization.get("segments", [])
        if not diar_segments:
            st.warning("Pyannote nie zwr√≥ci≈Ç poprawnych segment√≥w ‚Äî zwracam samƒÖ transkrypcjƒô.")
            return text_only, 1, meta

        output_lines = ["=== TRANSKRYPCJA Z IDENTYFIKACJƒÑ G≈ÅOS√ìW ===", ""]
        for seg in segments:
            start = float(seg.get("start", 0))
            end = float(seg.get("end", 0))
            txt = (seg.get("text") or "").strip()
            speaker = pick_speaker_for_interval(diar_segments, start, end)
            output_lines.append(f"[{start:.1f}s - {end:.1f}s] {speaker}: {txt}")

        result = "\n".join(output_lines)
        meta['has_speakers'] = True
        return result, 1, meta

    except Exception as e:
        logger.error(f"Audio with speakers error: {e}")
        return f"[B≈ÅƒÑD: {e}]", 0, {"type": "audio", "error": str(e)}


# === EKSTRAKTORY DOKUMENT√ìW ===
def extract_pdf(file, use_vision: bool, vision_model: str, ocr_pages_limit: int = 20):
    """PDF: tekst + opcjonalnie OCR/Vision (transkrypcja obrazu)."""
    texts = []
    try:
        file.seek(0)
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages):
                if i >= ocr_pages_limit:
                    texts.append(f"\n[... limit {ocr_pages_limit} stron ...]")
                    break
                t = page.extract_text() or ""
                texts.append(t)

        full_text = "\n".join(texts)

        if len(full_text.strip()) < MIN_TEXT_FOR_OCR_SKIP:
            file.seek(0)
            raw = file.read()
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                total_pages = len(pdf.pages)

            images = convert_from_bytes(
                raw,
                fmt="jpeg",
                dpi=150,
                first_page=1,
                last_page=min(ocr_pages_limit, total_pages)
            )

            if use_vision and vision_model:
                st.info(f"üñºÔ∏è U≈ºywam {vision_model} do analizy obraz√≥w...")
                for idx, img in enumerate(images[:ocr_pages_limit], 1):
                    st.caption(f"Przetwarzam stronƒô {idx}/{min(ocr_pages_limit, len(images))}")
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG", quality=85)
                    img_b64 = base64.b64encode(buf.getvalue()).decode()
                    response = query_ollama_vision(VISION_TRANSCRIBE_PROMPT, img_b64, vision_model)
                    texts.append(f"\n--- Strona {idx} (Vision) ---\n{response}")
            else:
                st.info("üìù OCR Tesseract...")
                for idx, img in enumerate(images[:ocr_pages_limit], 1):
                    st.caption(f"OCR strona {idx}/{min(ocr_pages_limit, len(images))}")
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG")
                    ocr_text = ocr_image_bytes(buf.getvalue())
                    texts.append(f"\n--- Strona {idx} (OCR) ---\n{ocr_text}")

        meta = {"type": "pdf", "pages": len(texts)}
        return "\n".join(texts), len(texts), meta
    except Exception as e:
        logger.error(f"PDF extract error: {e}")
        return f"[B≈ÅƒÑD PDF: {e}]", 0, {"type": "pdf", "error": str(e)}


def extract_pptx(file, use_vision: bool, vision_model: str):
    """PPTX: tekst + obrazy (opcjonalnie Vision opis)."""
    try:
        file.seek(0)
        prs = Presentation(file)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            parts = [f"=== Slajd {i} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    parts.append(f"Notatki: {notes}")

            if use_vision and vision_model:
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) == 13:  # PICTURE
                        try:
                            img_stream = shape.image.blob
                            img_b64 = base64.b64encode(img_stream).decode()
                            response = query_ollama_vision(VISION_DESCRIBE_PROMPT, img_b64, vision_model)
                            parts.append(f"[Obraz] {response}")
                        except Exception:
                            pass

            slides_text.append("\n".join(parts))

        return "\n\n".join(slides_text), len(prs.slides), {"type": "pptx", "slides": len(prs.slides)}
    except Exception as e:
        logger.error(f"PPTX error: {e}")
        return f"[B≈ÅƒÑD PPTX: {e}]", 0, {"type": "pptx", "error": str(e)}


def extract_docx(file):
    """DOCX: tekst + tabele."""
    try:
        file.seek(0)
        doc = Document(file)
        paras = [p.text for p in doc.paragraphs if p.text]

        for tbl in doc.tables:
            for row in tbl.rows:
                paras.append(" | ".join(cell.text for cell in row.cells))

        return "\n".join(paras), len(paras), {"type": "docx"}
    except Exception as e:
        logger.error(f"DOCX error: {e}")
        return f"[B≈ÅƒÑD DOCX: {e}]", 0, {"type": "docx", "error": str(e)}


# === ENHANCED VISION Z WEB SEARCH ===
def extract_keywords_from_vision(vision_text: str, max_keywords: int = 5) -> list:
    """WyciƒÖga kluczowe s≈Çowa z opisu Vision dla wyszukiwania."""
    stopwords = {
        'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
        'jest', 'sƒÖ', 'to', 'w', 'z', 'na', 'do', 'i', '≈ºe', 'siƒô', 'oraz'
    }

    words = re.findall(r'\b[a-zA-ZƒÖƒáƒô≈Ç≈Ñ√≥≈õ≈∫≈ºƒÑƒÜƒò≈Å≈É√ì≈ö≈π≈ª]{3,}\b', vision_text.lower())

    keywords = [w for w in words if w not in stopwords]
    unique = list(dict.fromkeys(keywords))
    sorted_by_length = sorted(unique, key=len, reverse=True)

    return sorted_by_length[:max_keywords]


def enhanced_vision_with_web_search(
    image_b64: str,
    prompt: str,
    vision_model: str,
    enable_web: bool = False
) -> str:
    """
    Vision + weryfikacja przez web search (opcjonalna).
    """
    logger.info(f"Vision analysis using {vision_model}...")
    vision_response = query_ollama_vision(prompt, image_b64, vision_model)

    if not enable_web or not st.session_state.get("ALLOW_WEB", False):
        return vision_response

    logger.info("Extracting keywords from vision response...")
    keywords = extract_keywords_from_vision(vision_response, max_keywords=3)

    if not keywords:
        logger.warning("No keywords extracted, returning basic vision response")
        return vision_response

    logger.info(f"Keywords for web search: {keywords}")

    try:
        search_query = " ".join(keywords)
        logger.info(f"Searching web for: '{search_query}'")

        web_results = web_search_and_summarize(
            queries=[search_query],
            max_results=2,
            model=st.session_state.get("selected_main_text_model", "qwen2.5:7b")
        )

        if not web_results.get("items"):
            logger.warning("No web results found, returning basic vision response")
            return vision_response

        web_context = "\n\n".join([
            f"≈πr√≥d≈Ço {i+1}: {item.get('title', 'N/A')}\n{item.get('summary', '')}"
            for i, item in enumerate(web_results["items"][:2])
        ])

        enhancement_prompt = f"""Masz dwa ≈∫r√≥d≈Ça informacji o obiekcie z obrazu:

ANALIZA OBRAZU (AI Vision):
{vision_response}

DODATKOWY KONTEKST Z INTERNETU:
{web_context}

ZADANIE:
Stw√≥rz OSTATECZNY, PRECYZYJNY opis obiektu z obrazu, u≈ºywajƒÖc:
1. Informacji z analizy obrazu (najwa≈ºniejsze - to co faktycznie widaƒá)
2. Kontekstu z internetu (uzupe≈Çnienie, weryfikacja terminologii)

ZASADY:
- Zachowaj strukturƒô 6-punktowƒÖ z opisu Vision
- Popraw b≈Çƒôdy terminologiczne je≈õli znajdziesz w kontek≈õcie internetowym
- Dodaj dodatkowe szczeg√≥≈Çy TYLKO je≈õli potwierdzajƒÖ to co widaƒá na obrazie
- Nie dodawaj informacji kt√≥rych nie ma na obrazie
- Pisz po polsku

Podaj ostateczny, ulepszony opis:"""

        logger.info("Enhancing vision response with web context...")
        enhanced = query_ollama_text(
            enhancement_prompt,
            model=st.session_state.get("selected_main_text_model"),
            json_mode=False,
            timeout=120
        )

        return enhanced

    except Exception as e:
        logger.error(f"Web search enhancement failed: {e}")
        return vision_response


def extract_image(file, use_vision: bool, vision_model: str, image_mode: str):
    """Obraz: OCR / Vision (przepisz) / Vision (opisz) / OCR+opis."""
    try:
        file.seek(0)
        img_bytes = file.read()
        results = []
        meta = {"type": "image", "mode": image_mode}

        if image_mode in ("ocr", "ocr_plus_vision_desc"):
            ocr_text = ocr_image_bytes(img_bytes)
            results.append(f"=== OCR ===\n{ocr_text}")

        if image_mode in ("vision_transcribe", "vision_describe", "ocr_plus_vision_desc"):
            if use_vision and vision_model:
                img_b64 = base64.b64encode(img_bytes).decode()
                prompt = VISION_TRANSCRIBE_PROMPT if image_mode == "vision_transcribe" else VISION_DESCRIBE_PROMPT

                enable_web_enhancement = st.session_state.get("ALLOW_WEB", False) and image_mode == "vision_describe"

                vis = enhanced_vision_with_web_search(
                    img_b64,
                    prompt,
                    vision_model,
                    enable_web=enable_web_enhancement
                )

                tag = "Vision (transkrypcja)" if image_mode == "vision_transcribe" else "Vision (opis)"
                results.append(f"=== {tag} ===\n{vis}")
                meta["vision_model"] = vision_model
                meta["web_enhanced"] = enable_web_enhancement
            else:
                results.append("[Vision niedostƒôpne]")

        txt = "\n\n".join(results).strip()
        return txt, 1, meta
    except Exception as e:
        logger.error(f"Image error: {e}")
        return f"[B≈ÅƒÑD IMG: {e}]", 0, {"type": "image", "error": str(e)}


# === E-MAIL: EML/MSG ===
def extract_eml(file):
    """EML: nag≈Ç√≥wki + tre≈õƒá. Wymaga mailparser (opcjonalnie fallback)."""
    try:
        file.seek(0)
        raw = file.read()
        if mailparser is None:
            text = raw.decode("utf-8", errors="ignore")
            return text, 1, {"type": "email", "note": "mailparser not installed"}
        m = mailparser.parse_from_bytes(raw)
        headers = [
            f"From: {m.from_[0][1] if m.from_ else ''}",
            f"To: {', '.join([x[1] for x in m.to]) if m.to else ''}",
            f"Subject: {m.subject or ''}",
            f"Date: {m.date.isoformat() if m.date else ''}"
        ]
        body = (m.text_plain[0] if m.text_plain else m.body) or ""
        attach = [att.get('filename') for att in (m.attachments or []) if att.get('filename')]
        if attach:
            headers.append(f"Attachments: {', '.join(attach)}")
        text = "\n".join(headers) + "\n\n" + (body or "")
        return text, 1, {"type": "email", "has_attachments": bool(attach)}
    except Exception as e:
        logger.error(f"EML error: {e}")
        return f"[B≈ÅƒÑD EML: {e}]", 0, {"type": "email", "error": str(e)}


def parse_msg_email(file):
    """MSG (Outlook): nag≈Ç√≥wki + tre≈õƒá. Wymaga extract_msg (fallback na surowy tekst)."""
    try:
        if extract_msg is None:
            file.seek(0)
            raw = file.read()
            text = raw.decode("utf-8", errors="ignore")
            return text, 1, {"type": "email", "note": "extract-msg not installed"}
        file.seek(0)
        raw = file.read()
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        msg = extract_msg.Message(tmp_path)
        headers = [
            f"From: {msg.sender or ''}",
            f"To: {msg.to or ''}",
            f"Subject: {msg.subject or ''}",
            f"Date: {msg.date or ''}"
        ]
        body = msg.body or ""
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        text = "\n".join(headers) + "\n\n" + body
        return text, 1, {"type": "email"}
    except Exception as e:
        logger.error(f"MSG error: {e}")
        return f"[B≈ÅƒÑD MSG: {e}]", 0, {"type": "email", "error": str(e)}


# === ROUTER ===
def process_file(file, use_vision: bool, vision_model: str, ocr_limit: int, image_mode: str,
                 enable_audio_chunking: bool = False, audio_chunk_minutes: int = 10):
    """
    Router do odpowiedniego ekstraktora.

    Args:
        file: Plik do przetworzenia
        use_vision: Czy u≈ºywaƒá vision models
        vision_model: Nazwa modelu vision
        ocr_limit: Limit stron dla OCR
        image_mode: Tryb przetwarzania obraz√≥w
        enable_audio_chunking: Czy w≈ÇƒÖczyƒá dzielenie audio na chunki
        audio_chunk_minutes: D≈Çugo≈õƒá chunka audio w minutach
    """
    name = file.name.lower()

    if name.endswith('.pdf'):
        return extract_pdf(file, use_vision, vision_model, ocr_limit)
    elif name.endswith(('.pptx', '.ppt')):
        return extract_pptx(file, use_vision, vision_model)
    elif name.endswith('.docx'):
        return extract_docx(file)
    elif name.endswith(('.jpg', '.jpeg', '.png')):
        return extract_image(file, use_vision, vision_model, image_mode)
    elif name.endswith(('.mp3', '.wav', '.m4a', '.ogg', '.flac')):
        ok, _ = check_pyannote_health(PYANNOTE_URL)
        if ok:
            return extract_audio_with_speakers(file, enable_audio_chunking, audio_chunk_minutes)
        return extract_audio_whisper(file, enable_audio_chunking, audio_chunk_minutes)
    elif name.endswith('.eml'):
        return extract_eml(file)
    elif name.endswith('.msg'):
        return parse_msg_email(file)
    elif name.endswith('.txt'):
        file.seek(0)
        content = file.read().decode('utf-8', errors='ignore')
        return content, 1, {"type": "txt"}
    else:
        return "[Nieobs≈Çugiwany format]", 0, {"type": "unknown"}


# === ANYTHINGLLM (lokalnie, respektuje offline) ===
def send_to_anythingllm(text: str, filename: str):
    """Wy≈õlij dokument do AnythingLLM (tylko gdy nie OFFLINE_MODE i URL prywatny)."""
    if OFFLINE_MODE:
        return False, "Tryb offline ‚Äî wysy≈Çka zablokowana"
    if not ANYTHINGLLM_URL or not ANYTHINGLLM_API_KEY:
        return False, "Brak konfiguracji AnythingLLM"
    try:
        headers = {"Authorization": f"Bearer {ANYTHINGLLM_API_KEY}"}
        payload = {"name": filename, "content": text, "type": "text/plain"}
        r = http_post(f"{ANYTHINGLLM_URL}/api/v1/document-upload", headers=headers, json=payload, timeout=30)
        r.raise_for_status()
        return True, "‚úÖ Wys≈Çano do AnythingLLM"
    except Exception as e:
        return False, f"B≈ÇƒÖd AnythingLLM: {e}"


# === PODSUMOWANIA AUDIO I DOKUMENT√ìW (Map-Reduce JSON + Markdown) ===
# === PODSUMOWANIA AUDIO I DOKUMENT√ìW (Map-Reduce JSON + Markdown) ===

MAP_PROMPT_TEMPLATE = """
Jeste≈õ asystentem ds. spotka≈Ñ (PL). Otrzymasz FRAGMENT transkrypcji rozmowy z klientem
(mogƒÖ wystƒôpowaƒá znaczniki typu SPEAKER_00, SPEAKER_01 oraz znaczniki czasu).

Twoje zadanie:
1. Zrozumieƒá, co faktycznie zosta≈Ço powiedziane w tym fragmencie.
2. WyciƒÖgnƒÖƒá najwa≈ºniejsze ustalenia, decyzje, zadania i ryzyka.

BARDZO WA≈ªNE ZASADY:
- Opieraj siƒô WY≈ÅƒÑCZNIE na tre≈õci fragmentu. Je≈õli czego≈õ w nim nie ma ‚Äì NIE wymy≈õlaj.
- Je≈õli nie masz informacji: u≈ºyj pustych list [] lub pustych string√≥w "".
- Nie dopisuj w≈Çasnych opinii, interpretacji ani ‚Äûco mog≈Çoby byƒá‚Äù.
- Je≈õli wystƒôpujƒÖ oznaczenia SPEAKER_x ‚Äì zachowaj je, nie t≈Çumacz ich na imiona.

WYMAGANY FORMAT WYJ≈öCIOWY (Tylko poprawny JSON, bez komentarzy, bez markdownu):
{{
  "summary": "1-2 kr√≥tkie akapity skr√≥tu (po polsku)",
  "key_points": ["punkt 1", "punkt 2"],
  "decisions": ["decyzja 1", "decyzja 2"],
  "to_be_decided": ["kwestia do ustalenia 1", "kwestia 2"],
  "action_items": [
    {{ "owner":"", "task":"", "due":"", "notes":"" }}
  ],
  "risks": [
    {{ "risk":"", "impact":"niski/≈õredni/wysoki", "mitigation":"" }}
  ],
  "open_questions": ["pytanie 1", "pytanie 2"]
}}

DODATKOWE WSKAZ√ìWKI:
- "summary" ma byƒá zwiƒôz≈Çym opisem tego fragmentu (maks. 5 zda≈Ñ).
- "key_points" to najwa≈ºniejsze merytoryczne informacje (fakty, ustalenia).
- "decisions" tylko wtedy, gdy wprost zapad≈Ça jaka≈õ decyzja.
- "to_be_decided" ‚Äì kwestie, kt√≥re WYNIKAJƒÑ z rozmowy jako wymagajƒÖce decyzji.
- "action_items" ‚Äì konkretne zadania; je≈õli da siƒô wskazaƒá w≈Ça≈õciciela (np. SPEAKER_01), u≈ºyj go w polu "owner".
- "risks" ‚Äì realne ryzyka, nie og√≥lne frazesy.

Je≈õli nie masz nic do wpisania w kt√≥re≈õ pole ‚Äì zwr√≥ƒá puste [] dla list lub "" dla tekstu.

FRAGMENT TRANSKRYPCJI:
{fragment}

Pamiƒôtaj: ODPOWIED≈π = TYLKO JEDEN OBIEKT JSON w powy≈ºszym formacie.
"""

REDUCE_PROMPT_TEMPLATE = """
Jeste≈õ asystentem ds. spotka≈Ñ (PL). Otrzymasz LISTƒò czƒô≈õciowych podsumowa≈Ñ w JSON
(z polami: summary, key_points, decisions, to_be_decided, action_items, risks, open_questions).

Twoje zadanie:
- Po≈ÇƒÖczyƒá te fragmenty w JEDNO sp√≥jne podsumowanie ca≈Çej rozmowy.
- UsunƒÖƒá duplikaty, scaliƒá podobne punkty, zachowaƒá sens wypowiedzi.

WEJ≈öCIE (lista JSON-√≥w z poprzedniego kroku):
{partials}

WYMAGANY FORMAT WYJ≈öCIOWY:
Zwr√≥ƒá TYLKO JEDEN obiekt JSON (bez komentarzy, bez markdownu), dok≈Çadnie w tej strukturze:
{{
  "summary": "skondensowany skr√≥t ca≈Ço≈õci (1-3 akapity, po polsku)",
  "key_points": ["..."],
  "decisions": ["..."],
  "to_be_decided": ["..."],
  "action_items": [
    {{ "owner":"", "task":"", "due":"", "notes":"" }}
  ],
  "risks": [
    {{ "risk":"", "impact":"niski/≈õredni/wysoki", "mitigation":"" }}
  ],
  "open_questions": ["..."]
}}

ZASADY:
- Usuwaj duplikaty w listach (ale nie usuwaj wa≈ºnych szczeg√≥≈Ç√≥w).
- Je≈õli jakie≈õ pole jest puste ‚Äì zwr√≥ƒá pustƒÖ listƒô [] lub pusty string "".
- Nie dodawaj nowych informacji, kt√≥rych nie ma w podsumowaniach wej≈õciowych.
- Pisz wy≈ÇƒÖcznie po polsku.

ODPOWIED≈π = TYLKO JSON, bez dodatkowego tekstu.
"""

# === PODSUMOWANIA DOKUMENT√ìW (PDF/DOCX/EMAIL/IMG) ===

DOC_SUMMARY_PROMPT = """
Jeste≈õ asystentem analizujƒÖcym dokumenty (PL). Otrzymasz TRE≈öƒÜ jednego dokumentu
(PDF/DOCX/PPTX/TXT/EMAIL lub opis po OCR/Vision).

Twoje zadanie:
- Stre≈õciƒá dokument.
- WyciƒÖgnƒÖƒá kluczowe informacje, potencjalne zadania, ryzyka i pytania.

FORMAT WYJ≈öCIOWY (TYLKO JSON, bez komentarzy/markdownu):
{{
  "summary": "1-3 akapity skr√≥tu (po polsku)",
  "key_points": ["punkt 1", "punkt 2"],
  "decisions": ["decyzja 1", "decyzja 2"],
  "action_items": [
    {{ "owner":"","task":"","due":"","priority":"low/medium/high","notes":"" }}
  ],
  "risks": [
    {{ "risk":"","impact":"low/medium/high","mitigation":"" }}
  ],
  "open_questions": ["pytanie 1","pytanie 2"]
}}

ZASADY:
- Nie wymy≈õlaj informacji ‚Äì je≈õli czego≈õ w tek≈õcie nie ma, zostaw puste [] lub "".
- "decisions" tylko, gdy wprost jest decyzja/ustalenie.
- "action_items" tylko przy konkretnych zadaniach (kto/co/kiedy), w razie braku owner/due zostaw "".
- Pisz wy≈ÇƒÖcznie po polsku.

TRE≈öƒÜ DOKUMENTU:
{content}
"""

DOC_SUMMARY_REDUCE_PROMPT = """
Jeste≈õ asystentem analizujƒÖcym dokumenty (PL). Otrzymasz listƒô czƒô≈õciowych podsumowa≈Ñ dokumentu w JSON.
Twoje zadanie to SCALIƒÜ je w jedno sp√≥jne podsumowanie w tym samym formacie JSON.

WEJ≈öCIE (lista JSON-√≥w):
{partials}

FORMAT WYJ≈öCIOWY (TYLKO JSON):
{{
  "summary": "1-3 akapity ca≈Ço≈õciowego skr√≥tu",
  "key_points": [...],
  "decisions": [...],
  "action_items": [...],
  "risks": [...],
  "open_questions": [...]
}}

ZASADY:
- Usu≈Ñ duplikaty, po≈ÇƒÖcz podobne punkty.
- Nie dodawaj informacji spoza wej≈õciowych JSON-√≥w.
- Pisz wy≈ÇƒÖcznie po polsku.
"""


def chunk_text(text: str, max_chars: int = 6000, overlap: int = 500) -> List[str]:
    if not text:
        return []
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + max_chars, n)
        chunks.append(text[start:end])
        if end >= n:
            break
        start = max(0, end - overlap)
    return chunks


def try_parse_json(s: str) -> Dict[str, Any]:
    if not s:
        return {}
    clean = s.strip()
    if clean.startswith("```json"):
        clean = clean[7:]
    if clean.startswith("```"):
        clean = clean[3:]
    if clean.endswith("```"):
        clean = clean[:-3]
    try:
        return json.loads(clean)
    except Exception:
        return {}


def merge_summary_dicts(items: List[Dict[str, Any]]) -> Dict[str, Any]:
    """≈ÅƒÖczy wiele s≈Çownik√≥w podsumowa≈Ñ w jeden (dla audio i dokument√≥w)."""
    out = {
        "summary": "",
        "key_points": [],
        "decisions": [],
        "to_be_decided": [],
        "action_items": [],
        "risks": [],
        "open_questions": []
    }
    for it in items:
        if not isinstance(it, dict):
            continue
        if it.get("summary"):
            out["summary"] += (("\n" if out["summary"] else "") + str(it["summary"]))
        for k in ["key_points", "decisions", "to_be_decided", "open_questions"]:
            vals = it.get(k, [])
            if isinstance(vals, list):
                out[k].extend(vals)
        for ai in it.get("action_items", []):
            if isinstance(ai, dict):
                out["action_items"].append(ai)
        for r in it.get("risks", []):
            if isinstance(r, dict):
                out["risks"].append(r)
    for k in ["key_points", "decisions", "to_be_decided", "open_questions"]:
        out[k] = list(dict.fromkeys(out[k]))
    return out


def build_meeting_summary_markdown(data: Dict[str, Any]) -> str:
    """Buduje markdown z JSON-owego podsumowania (u≈ºywane dla audio i dokument√≥w)."""
    if not data:
        return "_Brak danych do podsumowania_"
    md = []
    md.append("# Podsumowanie")

    if data.get("summary"):
        md.append(data["summary"])

    md.append("\n## Kluczowe punkty")
    for x in (data.get("key_points") or []) or ["brak"]:
        md.append(f"- {x}")

    md.append("\n## Decyzje vs Do ustalenia")
    md.append("### Decyzje")
    for d in (data.get("decisions") or []) or ["brak"]:
        md.append(f"- {d}")
    md.append("### Do ustalenia")
    for q in (data.get("to_be_decided") or []) or ["brak"]:
        md.append(f"- {q}")

    md.append("\n## Zadania (Action Items)")
    action_items = data.get("action_items", [])
    if not action_items:
        md.append("- brak")
    else:
        for ai in action_items:
            # Handle both dict and string formats
            if isinstance(ai, dict):
                owner = ai.get("owner", "") or "N/A"
                task = ai.get("task", "") or "N/A"
                due = ai.get("due", "") or "-"
                notes = ai.get("notes", "") or ""
                md.append(f"- [ ] {task} (owner: {owner}, termin: {due}) {('- ' + notes) if notes else ''}")
            elif isinstance(ai, str):
                md.append(f"- [ ] {ai}")
            else:
                md.append(f"- [ ] {str(ai)}")

    md.append("\n## Ryzyka")
    risks = data.get("risks", [])
    if not risks:
        md.append("- brak")
    else:
        for r in risks:
            # Handle both dict and string formats
            if isinstance(r, dict):
                risk = r.get("risk", "")
                impact = r.get("impact", "")
                mit = r.get("mitigation", "")
                md.append(f"- {risk} (wp≈Çyw: {impact}) ‚Üí mitygacja: {mit}")
            elif isinstance(r, str):
                md.append(f"- {r}")
            else:
                md.append(f"- {str(r)}")

    md.append("\n## Pytania / Otwarte kwestie")
    for q in (data.get("open_questions") or []) or ["brak"]:
        md.append(f"- {q}")

    return "\n".join(md)


def summarize_meeting_transcript(
    transcript: str,
    model: str = "llama3:latest",
    max_chars: int = 6000,
    diarized: bool = False
) -> Dict[str, Any]:
    """Map-Reduce podsumowanie transkrypcji rozmowy (audio)."""
    if not transcript or len(transcript.strip()) < 20:
        return {}

    parts = chunk_text(transcript, max_chars=max_chars, overlap=500)
    partials: List[Dict[str, Any]] = []

    for p in parts:
        prompt = MAP_PROMPT_TEMPLATE.format(fragment=p)
        resp = query_ollama_text(prompt, model=model, json_mode=True, timeout=180)
        data = try_parse_json(resp)
        if not data:
            resp2 = query_ollama_text(prompt, model=model, json_mode=False, timeout=180)
            data = try_parse_json(resp2)
        if data:
            partials.append(data)

    if not partials:
        return {"summary": transcript[:1200] + ("..." if len(transcript) > 1200 else "")}

    partials_str = json.dumps(partials, ensure_ascii=False, indent=2)
    reduce_prompt = REDUCE_PROMPT_TEMPLATE.format(partials=partials_str)
    reduce_resp = query_ollama_text(reduce_prompt, model=model, json_mode=True, timeout=240)
    final_data = try_parse_json(reduce_resp)
    if not final_data:
        final_data = merge_summary_dicts(partials)
    return final_data


def summarize_document_text(
    text: str,
    model: str,
    max_chars: int = 6000
) -> Dict[str, Any]:
    """Og√≥lne podsumowanie dokumentu tekstowego (PDF/DOCX/PPTX/TXT/EMAIL/IMG-OCR)."""
    if not text or len(text.strip()) < 20:
        return {}

    parts = chunk_text(text, max_chars=max_chars, overlap=500)
    partials: List[Dict[str, Any]] = []

    for p in parts:
        prompt = DOC_SUMMARY_PROMPT.format(content=p)
        resp = query_ollama_text(prompt, model=model, json_mode=True, timeout=180)
        data = try_parse_json(resp)
        if not data:
            resp2 = query_ollama_text(prompt, model=model, json_mode=False, timeout=180)
            data = try_parse_json(resp2)
        if data:
            partials.append(data)

    if not partials:
        return {"summary": text[:1200] + ("..." if len(text) > 1200 else "")}

    partials_str = json.dumps(partials, ensure_ascii=False, indent=2)
    reduce_prompt = DOC_SUMMARY_REDUCE_PROMPT.format(partials=partials_str)
    reduce_resp = query_ollama_text(reduce_prompt, model=model, json_mode=True, timeout=240)
    final_data = try_parse_json(reduce_resp)
    if not final_data:
        final_data = merge_summary_dicts(partials)
    return final_data


# === PROJECT BRAIN (klasyfikacja, zadania, ryzyka, brief, web queries) ===
# === PROJECT BRAIN (klasyfikacja, zadania, ryzyka, brief, web queries) ===

DOC_CLASS_PROMPT = """
Jeste≈õ asystentem analizujƒÖcym dokumenty (PL).

Zadanie:
Na podstawie PODANEJ TRE≈öCI okre≈õl, jaki to typ dokumentu.

Dozwolone typy:
["email","chat","meeting_transcript","spec","invoice","drawing","image_text","note","other"]

FORMAT WYJ≈öCIA (TYLKO JSON):
{{
  "type": "<jeden_z_powy≈ºszych_typ√≥w>"
}}

TRE≈öƒÜ (poczƒÖtek dokumentu):
{content}
"""

TASKS_PROMPT = """
Jeste≈õ asystentem PM (PL). Z tre≈õci dokumentu wyodrƒôbnij listƒô KONKRETNYCH zada≈Ñ do wykonania.

FORMAT WYJ≈öCIA (TYLKO JSON):
{{
 "tasks":[
   {{ "owner":"","task":"","due":"","priority":"low/medium/high","tags":[],"source":"" }}
 ]
}}

ZASADY:
- "task" musi byƒá konkretnym dzia≈Çaniem (co dok≈Çadnie zrobiƒá).
- "owner" je≈õli mo≈ºliwy, inaczej "".
- "due" je≈õli jest termin (data, tydzie≈Ñ, "ASAP"), inaczej "".
- "priority" oszacuj na podstawie tre≈õci: low/medium/high.
- "tags" ‚Äì lista kr√≥tkich tag√≥w (np. ["CAD","dokumentacja"]), je≈õli nie ma ‚Üí [].
- "source" ‚Äì kr√≥tki opis skƒÖd w dokumencie wynika to zadanie (np. "sekcja wymagania").

TRE≈öƒÜ:
{content}
"""

RISKS_PROMPT = """
Jeste≈õ asystentem PM (PL). Z tre≈õci dokumentu wyodrƒôbnij:
- ryzyka,
- za≈Ço≈ºenia (assumptions),
- pytania do wyja≈õnienia (RFI).

FORMAT WYJ≈öCIA (TYLKO JSON):
{{
 "risks":[{{"risk":"","impact":"low/medium/high","mitigation":""}}],
 "assumptions":["..."],
 "rfis":["pytanie 1","pytanie 2"]
}}

ZASADY:
- "risks" ‚Äì tylko realne ryzyka wynikajƒÖce z tre≈õci (techniczne, zakresowe, harmonogramowe, organizacyjne).
- "impact" ‚Äì oszacowany wp≈Çyw: low/medium/high.
- "mitigation" ‚Äì realna propozycja dzia≈Çania zmniejszajƒÖcego ryzyko.
- "assumptions" ‚Äì rzeczy przyjƒôte jako oczywiste w dokumencie.
- "rfis" ‚Äì pytania, kt√≥re warto zadaƒá klientowi.

TRE≈öƒÜ:
{content}
"""

PROJECT_BRIEF_PROMPT = """
Masz wiele FRAGMENT√ìW projektu (podsumowania, zadania, ryzyka).

Wej≈õcie (JSON LISTA element√≥w):
{items}

Ka≈ºdy element mo≈ºe zawieraƒá m.in.:
- "summary"
- "key_points"
- "decisions"
- "risks"
- "action_items"
- "rfis"

Twoje zadanie: Stw√≥rz projektowy BRIEF dla PM.

FORMAT WYJ≈öCIA (TYLKO JSON):
{{
 "brief":"1-3 akapity podsumowania projektu (PL)",
 "key_points":[],
 "decisions":[],
 "rfis":[],
 "risks":[{{"risk":"","impact":"low/medium/high","mitigation":""}}],
 "next_steps":[]
}}

ZASADY:
- "brief" ma byƒá zrozumia≈Çy dla PM kt√≥ry nie czyta≈Ç dokumentu.
- "key_points" ‚Äì najwa≈ºniejsze fakty.
- "decisions" ‚Äì tylko decyzje kt√≥re faktycznie wynikajƒÖ z wej≈õcia.
- "rfis" ‚Äì pytania do klienta/zamawiajƒÖcego.
- "next_steps" ‚Äì konkretne nastƒôpne kroki (co zrobiƒá dalej).

Zwr√≥ƒá WY≈ÅƒÑCZNIE JSON w powy≈ºszym formacie.
"""

WEB_QUERIES_PROMPT = """
Na bazie tre≈õci dokumentu zaproponuj 3‚Äì5 neutralnych zapyta≈Ñ do wyszukiwarki,
bez danych wra≈ºliwych (bez nazw firm, os√≥b, maili, numer√≥w).

FORMAT WYJ≈öCIA (TYLKO JSON):
{{"queries":["zapytanie 1","zapytanie 2"]}}

TRE≈öƒÜ:
{content}
"""


def classify_document(text: str, model: str) -> str:
    resp = query_ollama_text(DOC_CLASS_PROMPT.format(content=text[:4000]), model=model, json_mode=True, timeout=90)
    data = try_parse_json(resp)
    return (data.get("type") or "other") if isinstance(data, dict) else "other"


def extract_tasks_from_text(text: str, model: str) -> List[Dict[str, Any]]:
    resp = query_ollama_text(TASKS_PROMPT.format(content=text[:6000]), model=model, json_mode=True, timeout=120)
    data = try_parse_json(resp) or {}
    tasks = data.get("tasks", []) if isinstance(data, dict) else []
    out = []
    for t in tasks:
        if isinstance(t, dict):
            out.append({
                "owner": t.get("owner", ""),
                "task": t.get("task", ""),
                "due": t.get("due", ""),
                "priority": t.get("priority", ""),
                "tags": t.get("tags", []),
                "source": t.get("source", "")
            })
    return out


def extract_risks_from_text(text: str, model: str) -> Dict[str, Any]:
    resp = query_ollama_text(RISKS_PROMPT.format(content=text[:6000]), model=model, json_mode=True, timeout=120)
    data = try_parse_json(resp) or {}
    return {
        "risks": data.get("risks", []),
        "assumptions": data.get("assumptions", []),
        "rfis": data.get("rfis", [])
    }


def build_project_brief(items: List[Dict[str, Any]], model: str) -> Dict[str, Any]:
    payload = json.dumps(items, ensure_ascii=False)[:12000]
    resp = query_ollama_text(PROJECT_BRIEF_PROMPT.format(items=payload), model=model, json_mode=True, timeout=180)
    data = try_parse_json(resp)
    if not data:
        return {"brief": "", "key_points": [], "decisions": [], "rfis": [], "risks": [], "next_steps": []}
    return data


def propose_web_queries(text: str, model: str) -> List[str]:
    resp = query_ollama_text(WEB_QUERIES_PROMPT.format(content=text[:4000]), model=model, json_mode=True, timeout=90)
    data = try_parse_json(resp) or {}
    q = data.get("queries", [])
    clean = []
    for s in q:
        s = re.sub(r'\S+@\S+', '[email]', s)
        s = re.sub(r'\b\d{7,}\b', '[num]', s)
        clean.append(s)
    return clean


def web_search_and_summarize(queries: List[str], max_results: int, model: str) -> Dict[str, Any]:
    """Pobiera strony przez DuckDuckGo i streszcza lokalnie.
       Dzia≈Ça tylko, gdy w UI w≈ÇƒÖczono ALLOW_WEB (sesja)."""
    if not st.session_state.get("ALLOW_WEB", False):
        return {"note": "web lookup disabled (ALLOW_WEB=False)", "items": []}
    if DDGS is None or trafilatura is None:
        return {"note": "duckduckgo-search/trafilatura not installed", "items": []}
    results = []
    try:
        with DDGS() as ddg:
            for q in queries or []:
                hits = ddg.text(q, region="pl-pl", safesearch="moderate", max_results=max_results)
                for h in (hits or []):
                    url = h.get("href") or h.get("url")
                    title = h.get("title", "")
                    if not url:
                        continue
                    try:
                        c = trafilatura.fetch_url(url)
                        txt = trafilatura.extract(c) or ""
                    except Exception:
                        txt = ""
                    if not txt:
                        continue
                    summary = query_ollama_text(
                        f"Stre≈õƒá (PL) w 5 punktach:\n\n{txt[:6000]}",
                        model=model,
                        json_mode=False,
                        timeout=90
                    )
                    results.append({"query": q, "url": url, "title": title, "summary": summary})
                    if len(results) >= max_results:
                        break
                if len(results) >= max_results:
                    break
    except Exception as e:
        logger.error(f"Web lookup error: {e}")
    return {"items": results}


# Czƒô≈õƒá 3/3: UI ‚Äî sidebar, uploader, konwersja, zapisy, podsumowania audio+docs,
#            Project Brain (zadania/ryzyka/brief) + opcjonalny web lookup

# === FLAGI BIEGU / STOP ===
def init_run_flags():
    ss = st.session_state
    ss.setdefault("converting", False)
    ss.setdefault("cancel_requested", False)


def start_conversion():
    st.session_state["converting"] = True
    st.session_state["cancel_requested"] = False


def request_cancel():
    st.session_state["cancel_requested"] = True


def end_conversion():
    st.session_state["converting"] = False


init_run_flags()

# === UI / SIDEBAR ===
st.title("üìÑ Document Converter Pro")
st.caption("Konwersja PDF/DOCX/PPTX/IMG/AUDIO/EMAIL ‚Üí TXT z OCR, Vision lub Whisper (offline-first)")

with st.sidebar:
    st.header("‚öôÔ∏è Ustawienia")

    # === MODELE AI ===
    st.subheader("ü§ñ Modele AI")

    with st.expander("‚ÑπÔ∏è Co to sƒÖ modele AI?", expanded=False):
        st.markdown("""
        **Modele AI** to "m√≥zgi" aplikacji, kt√≥re przetwarzajƒÖ tekst i obrazy.

        - **Model tekstowy** - analizuje i podsumowuje dokumenty
        - **Model wizyjny** - rozpoznaje tre≈õƒá na obrazach/zdjƒôciach

        üí° **Porada:** Wiƒôksze modele (np. 14b) sƒÖ dok≈Çadniejsze ale wolniejsze.
        Mniejsze (7b) sƒÖ szybsze ale mniej dok≈Çadne.
        """)

    # Model tekstowy (g≈Ç√≥wny)
    available_text_models = [
        m for m in list_ollama_models()
        if not any(m.startswith(p) for p in ("llava", "bakllava", "moondream", "qwen2-vl", "qwen2.5vl", "nomic-embed"))
    ]

    if "selected_main_text_model" not in st.session_state:
        default_text = "qwen2.5:14b" if "qwen2.5:14b" in available_text_models else (
            available_text_models[0] if available_text_models else "llama3:latest"
        )
        st.session_state["selected_main_text_model"] = default_text

    try:
        text_idx = available_text_models.index(st.session_state["selected_main_text_model"])
    except (ValueError, IndexError):
        text_idx = 0
        st.session_state["selected_main_text_model"] = (
            available_text_models[0] if available_text_models else "llama3:latest"
        )

    main_text_model = st.selectbox(
        "Model tekstowy (g≈Ç√≥wny)",
        options=available_text_models or ["llama3:latest"],
        index=text_idx,
        key="main_text_sel",
        help="U≈ºywany do: web search enhancement, Project Brain, wszystkich operacji tekstowych",
        disabled=st.session_state.get("converting", False)
    )
    st.session_state["selected_main_text_model"] = main_text_model

    st.markdown("---")

    # === PRYWATNO≈öƒÜ I INTERNET ===
    st.subheader("üîí Prywatno≈õƒá i Internet")

    with st.expander("‚ÑπÔ∏è Co to znaczy?", expanded=False):
        st.markdown("""
        **Tryb offline** - blokuje wszystkie po≈ÇƒÖczenia internetowe poza lokalnymi us≈Çugami.

        **Web lookup** - pozwala aplikacji pobieraƒá publiczne strony WWW dla uzupe≈Çnienia informacji.

        ‚ö†Ô∏è **WA≈ªNE:** Aplikacja NIE wysy≈Ça Twoich dokument√≥w na zewnƒÖtrz!
        Web lookup pobiera TYLKO publiczne strony (np. Wikipedia) jako kontekst.
        """)

    OFFLINE_MODE = st.checkbox(
        "üîê Tryb offline (maksymalna prywatno≈õƒá)",
        value=OFFLINE_MODE,
        help="Blokuje dostƒôp do internetu. U≈ºywa tylko lokalnych us≈Çug.",
        disabled=st.session_state.get("converting", False)
    )
    st.session_state["ALLOW_WEB"] = st.checkbox(
        "üåê Web lookup (pobieranie publicznych stron)",
        value=st.session_state.get("ALLOW_WEB", True),
        help="Pozwala pobieraƒá publiczne strony WWW (bez wysy≈Çania Twoich dokument√≥w).",
        disabled=st.session_state.get("converting", False) or OFFLINE_MODE
    )

    if st.session_state.get("ALLOW_WEB", False):
        st.info("üîç Web search aktywny - Vision/Project Brain mogƒÖ u≈ºywaƒá kontekstu z sieci")
    else:
        st.success("üîí Web search wy≈ÇƒÖczony - maksymalna prywatno≈õƒá")

    with st.expander("üîå Status us≈Çug", expanded=False):
        st.caption("üìä Po≈ÇƒÖczenia z lokalnymi us≈Çugami:")

        def _status_url(name, url, desc=""):
            try:
                host = urlparse(url).hostname or ""
                is_local = is_private_host(host)
                icon = "‚úÖ" if is_local else "‚ùå"
                status = "lokalny" if is_local else "zewnƒôtrzny"
                st.caption(f"{icon} **{name}** {desc}")
                st.caption(f"   ‚îî‚îÄ `{url}` ({status})")
            except Exception:
                st.caption(f"‚ö†Ô∏è **{name}** - nie mo≈ºna zweryfikowaƒá")

        _status_url("Ollama", OLLAMA_URL, "- AI models")
        _status_url("Whisper", WHISPER_URL, "- Transkrypcja audio")
        _status_url("Pyannote", PYANNOTE_URL, "- Rozpoznawanie m√≥wc√≥w")

    st.markdown("---")

    # === VISION ===
    st.subheader("üëÅÔ∏è Vision (analiza obraz√≥w)")

    with st.expander("‚ÑπÔ∏è Co to jest Vision?", expanded=False):
        st.markdown("""
        **Vision** to AI kt√≥ry "widzi" obrazy i potrafi je opisaƒá lub przeczytaƒá tekst z nich.

        **Tryby pracy:**
        - **OCR** - tylko rozpoznawanie tekstu (Tesseract)
        - **Vision: przepisz tekst** - AI czyta tekst z obrazu (lepsze od OCR)
        - **Vision: opisz obraz** - AI opisuje CO WIDZI na obrazie
        - **OCR + Vision** - oba razem
        """)

    vision_models = list_vision_models()
    use_vision = st.checkbox(
        "‚ú® W≈ÇƒÖcz Vision (AI dla obraz√≥w)",
        value=True if vision_models else False,
        help="U≈ºywa AI do analizy obraz√≥w, zdjƒôƒá, schemat√≥w, rysunk√≥w technicznych",
        disabled=st.session_state.get("converting", False)
    )

    if vision_models and use_vision:
        if "selected_vision_model" not in st.session_state:
            default_vision = "qwen2.5vl:7b"
            st.session_state["selected_vision_model"] = (
                default_vision if default_vision in vision_models else
                next((m for m in vision_models if m.startswith("qwen")), vision_models[0])
            )

        try:
            vision_idx = vision_models.index(st.session_state["selected_vision_model"])
        except (ValueError, IndexError):
            vision_idx = 0
            st.session_state["selected_vision_model"] = vision_models[0]

        selected_vision = st.selectbox(
            "Model wizyjny (obrazy/rysunki)",
            vision_models,
            index=vision_idx,
            key="vision_model_sel",
            help="Model do analizy zdjƒôƒá, schemat√≥w, rysunk√≥w technicznych",
            disabled=st.session_state.get("converting", False)
        )
        st.session_state["selected_vision_model"] = selected_vision
    else:
        selected_vision = None
        if use_vision:
            st.warning("‚ö†Ô∏è Brak modeli Vision w Ollama (np. llava:13b / qwen2-vl:7b)")

    st.subheader("OCR")
    ocr_pages_limit = st.slider(
        "Limit stron OCR", 5, 50, 20,
        disabled=st.session_state.get("converting", False)
    )

    st.markdown("---")

    # === OPCJE ZAAWANSOWANE ===
    with st.expander("üîß Opcje zaawansowane", expanded=False):
        st.subheader("Obrazy (IMG)")

    if use_vision and selected_vision:
        if "image_mode_idx" not in st.session_state:
            st.session_state["image_mode_idx"] = 2  # "Vision: opisz obraz"

        image_mode_label = st.selectbox(
            "Tryb dla obraz√≥w",
            options=["OCR", "Vision: przepisz tekst", "Vision: opisz obraz", "OCR + Vision opis"],
            index=st.session_state["image_mode_idx"],
            key="img_mode_sel",
            disabled=st.session_state.get("converting", False)
        )
        st.session_state["image_mode_idx"] = [
            "OCR", "Vision: przepisz tekst", "Vision: opisz obraz", "OCR + Vision opis"
        ].index(image_mode_label)
    else:
        image_mode_label = st.selectbox(
            "Tryb dla obraz√≥w",
            options=["OCR"],
            index=0,
            disabled=True
        )
    image_mode = IMAGE_MODE_MAP.get(image_mode_label, "ocr")

    # === AUDIO PROCESSING ===
    st.markdown("---")
    st.subheader("üéôÔ∏è Przetwarzanie Audio")

    with st.expander("‚ÑπÔ∏è Po co dzielenie audio na czƒô≈õci?", expanded=False):
        st.markdown("""
        **Audio Chunking** dzieli d≈Çugie nagrania (60+ min) na mniejsze czƒô≈õci przed wys≈Çaniem do Whisper/Pyannote.

        **Korzy≈õci:**
        - ‚úÖ Unika timeout√≥w dla d≈Çugich plik√≥w
        - ‚úÖ Mniejsze zu≈ºycie pamiƒôci
        - ‚úÖ Lepsza niezawodno≈õƒá przetwarzania

        **Kiedy w≈ÇƒÖczyƒá?**
        - Nagrania > 30-60 minut
        - Problemy z timeoutami
        - Ograniczona pamiƒôƒá serwera
        """)

    enable_audio_chunking = st.checkbox(
        "üî™ Dziel d≈Çugie audio na czƒô≈õci (chunking)",
        value=False,
        help="W≈ÇƒÖcz dla nagra≈Ñ >60min lub gdy wystƒôpujƒÖ timeouty",
        disabled=st.session_state.get("converting", False)
    )

    audio_chunk_minutes = st.slider(
        "D≈Çugo≈õƒá chunka audio (minuty)",
        min_value=5,
        max_value=30,
        value=10,
        step=5,
        help="Kr√≥tsze chunki = bezpieczniej, d≈Çu≈ºsze = szybciej",
        disabled=st.session_state.get("converting", False) or not enable_audio_chunking
    )

    st.session_state["enable_audio_chunking"] = enable_audio_chunking
    st.session_state["audio_chunk_minutes"] = audio_chunk_minutes

    if enable_audio_chunking:
        st.info(f"‚úÇÔ∏è Audio bƒôdzie dzielone na czƒô≈õci po {audio_chunk_minutes} min")

    # Zapis lokalny
    st.markdown("---")
    st.subheader("Zapis lokalny")
    enable_local_save = st.checkbox(
        "Zapisz wyniki lokalnie (folder)", value=False,
        disabled=st.session_state.get("converting", False)
    )
    base_output_dir = st.text_input(
        "Katalog wyj≈õciowy", value="outputs",
        disabled=st.session_state.get("converting", False)
    )

    # AnythingLLM
    st.subheader("AnythingLLM")
    has_anythingllm_cfg = bool(ANYTHINGLLM_URL and ANYTHINGLLM_API_KEY)
    if OFFLINE_MODE:
        st.caption("Status: üîí Wy≈ÇƒÖczone (tryb offline)")
        has_anythingllm = False
    else:
        st.caption(f"Status: {'‚úÖ Skonfigurowane' if has_anythingllm_cfg else '‚ùå Brak config'}")
        has_anythingllm = has_anythingllm_cfg

    # === PODSUMOWANIA AI (audio + dokumenty) ===
    st.subheader("üß† Podsumowania AI")

    summarize_audio_enabled = st.checkbox(
        "Podsumowanie rozm√≥w audio",
        value=True,
        disabled=st.session_state.get("converting", False),
    )
    summarize_text_enabled = st.checkbox(
        "Podsumowanie dokument√≥w tekstowych (PDF/DOCX/PPTX/TXT/EMAIL)",
        value=False,
        disabled=st.session_state.get("converting", False),
    )
    summarize_images_enabled = st.checkbox(
        "Podsumowanie opis√≥w z obraz√≥w/rysunk√≥w",
        value=False,
        disabled=st.session_state.get("converting", False),
    )

    st.session_state["summarize_audio_enabled"] = summarize_audio_enabled
    st.session_state["summarize_text_enabled"] = summarize_text_enabled
    st.session_state["summarize_images_enabled"] = summarize_images_enabled

    summarize_model_candidates = [
        m for m in list_ollama_models()
        if not any(m.startswith(p) for p in ("llava", "bakllava", "moondream", "llava-phi", "nomic-embed", "qwen2-vl"))
    ]

    if "selected_summary_model" not in st.session_state:
        default = "qwen2.5:7b"
        st.session_state["selected_summary_model"] = (
            default if default in summarize_model_candidates else
            (summarize_model_candidates[0] if summarize_model_candidates else "llama3:latest")
        )

    try:
        sum_idx = summarize_model_candidates.index(st.session_state["selected_summary_model"])
    except (ValueError, IndexError):
        sum_idx = 0

    summarize_model = st.selectbox(
        "Model do podsumowa≈Ñ",
        options=summarize_model_candidates or ["llama3:latest"],
        index=sum_idx,
        key="sum_model_sel",
        disabled=st.session_state.get("converting", False)
    )
    st.session_state["selected_summary_model"] = summarize_model

    chunk_chars = st.slider(
        "Rozmiar chunku (znaki)", min_value=2000, max_value=8000, value=6000, step=500,
        disabled=st.session_state.get("converting", False)
    )

    # Project Brain toggle (mo≈ºesz wykorzystaƒá dalej w UI)
    st.subheader("üß≠ Project Brain (PM)")
    enable_project_brain = st.checkbox(
        "W≈ÇƒÖcz Project Brain (zadania/ryzyka/brief)", value=True,
        disabled=st.session_state.get("converting", False)
    )

    # Diagnostyka
    st.subheader("üîß Diagnostyka ≈õrodowiska")
    if st.button("Skanuj ≈õrodowisko", disabled=st.session_state.get("converting", False)):
        st.session_state["diag"] = run_diagnostics()
    if st.session_state.get("diag"):
        diag = st.session_state["diag"]
        st.caption("Wyniki diagnostyki (skr√≥t):")
        miss = diag.get("missing", {})
        st.write(f"- BrakujƒÖce systemowe: {', '.join(miss.get('system_packages', []) ) or '‚Äî'}")
        st.write(f"- BrakujƒÖce jƒôzyki Tesseract: {', '.join(miss.get('tesseract_languages', []) ) or '‚Äî'}")
        st.write(f"- BrakujƒÖce modu≈Çy Python: {', '.join(miss.get('python_modules', []) ) or '‚Äî'}")
        rec = diag.get("recommend_install", {})
        if rec.get("zypper") or rec.get("pip"):
            st.markdown("Zalecane instalacje:")
            if rec.get("zypper"):
                st.code("sudo zypper in -y " + " ".join(sorted(set(rec["zypper"]))), language="bash")
            if rec.get("pip"):
                st.code("pip install " + " ".join(sorted(set(rec["pip"]))), language="bash")
        with st.expander("Pe≈Çne szczeg√≥≈Çy diagnostyki"):
            st.json(diag, expanded=False)

    st.markdown("---")
    with st.expander("‚ùì Pomoc i podpowiedzi", expanded=False):
        st.markdown("""
        ### üéØ Szybki start

        1. **Upload pliku** - PDF, Word, zdjƒôcie, audio
        2. **Kliknij "Konwertuj"**
        3. **Gotowe!**

        **Dla audio:**
        - Automatycznie u≈ºywa Whisper (transkrypcja)
        - Pyannote rozpoznaje m√≥wc√≥w (je≈õli dostƒôpny)
        """)

# === FILE UPLOADER ===
uploaded_files = st.file_uploader(
    "Wgraj dokumenty",
    type=['pdf', 'docx', 'pptx', 'ppt', 'jpg', 'jpeg', 'png', 'txt', 'mp3', 'wav', 'm4a', 'ogg', 'flac'],
    accept_multiple_files=True,
    disabled=st.session_state.get("converting", False)
)

# STOP podczas konwersji
if st.session_state.get("converting", False):
    st.info("‚è≥ Trwa konwersja. Pozosta≈Çe akcje sƒÖ zablokowane.")
    if st.button("üü• STOP konwersji", type="primary"):
        request_cancel()

# === KONWERSJA ‚Üí zapis do session_state ===
if uploaded_files:
    st.info(f"üìÅ {len(uploaded_files)} plik√≥w")

    if st.button("üöÄ Konwertuj wszystkie", type="primary", key="btn_convert_all",
                 disabled=st.session_state.get("converting", False)):
        start_conversion()
        st.rerun()

# PƒòTLA KONWERSJI (wykonuje siƒô gdy converting=True)
if st.session_state.get("converting", False):
    if "conversion_started" not in st.session_state:
        st.session_state["conversion_started"] = True
        st.session_state["results"] = []
        st.session_state["combined_text"] = ""
        st.session_state["audio_items"] = []
        st.session_state["audio_summaries"] = []
        st.session_state["doc_summaries"] = []
        st.session_state["stats"] = {'processed': 0, 'errors': 0, 'pages': 0}
        st.session_state["converted"] = False
        st.session_state["files_sig"] = files_signature(uploaded_files)
        st.session_state["run_dir"] = create_run_dir(base_output_dir) if enable_local_save else None
        if enable_local_save:
            st.info(f"üíæ Wyniki bƒôdƒÖ zapisane w: {st.session_state['run_dir']}")

    progress = st.progress(0)
    all_texts = []

    for idx, file in enumerate(uploaded_files):
        if st.session_state.get("cancel_requested"):
            st.warning("‚õî Przerwano na ≈ºƒÖdanie u≈ºytkownika.")
            end_conversion()
            if "conversion_started" in st.session_state:
                del st.session_state["conversion_started"]
            st.rerun()

        try:
            progress.progress((idx + 1) / len(uploaded_files), text=f"Przetwarzam: {file.name}")
        except TypeError:
            progress.progress((idx + 1) / len(uploaded_files))

        st.subheader(f"üìÑ {file.name}")

        try:
            # Pobierz ustawienia audio chunking z session state
            enable_audio_chunking = st.session_state.get("enable_audio_chunking", False)
            audio_chunk_minutes = st.session_state.get("audio_chunk_minutes", 10)

            extracted_text, pages, meta = process_file(
                file, use_vision, selected_vision, ocr_pages_limit, image_mode,
                enable_audio_chunking, audio_chunk_minutes
            )

            st.session_state["results"].append({
                "name": file.name,
                "text": extracted_text,
                "original_text": extracted_text,
                "meta": meta,
                "pages": pages
            })

            all_texts.append(f"\n{'=' * 80}\n")
            all_texts.append(f"PLIK: {file.name}\n")
            all_texts.append(
                f"Typ: {getattr(file, 'type', 'unknown')}, Rozmiar: {getattr(file, 'size', 0) / 1024:.1f} KB\n"
            )
            all_texts.append(f"{'=' * 80}\n")
            all_texts.append(extracted_text)
            all_texts.append(f"\n[Stron/sekcji: {pages}]\n")

            st.session_state["stats"]["processed"] += 1
            st.session_state["stats"]["pages"] += pages

            with st.expander(f"Preview: {file.name}"):
                st.text(extracted_text[:2000] + ("..." if len(extracted_text) > 2000 else ""))

            if isinstance(meta, dict) and meta.get("type") == "audio":
                st.session_state["audio_items"].append((file.name, extracted_text, meta))

        except Exception as e:
            st.error(f"‚ùå B≈ÇƒÖd: {e}")
            logger.exception(f"Error processing {file.name}")
            st.session_state["stats"]["errors"] += 1

    progress.empty()
    st.session_state["combined_text"] = "\n".join(all_texts)
    st.session_state["converted"] = True
    end_conversion()
    if "conversion_started" in st.session_state:
        del st.session_state["conversion_started"]
    st.rerun()

# === SEKCJA WYNIK√ìW ===
if st.session_state.get("converted"):
    st.success(
        f"‚úÖ Przetworzono: {st.session_state['stats']['processed']} plik√≥w | Sekcji: {st.session_state['stats']['pages']}"
    )

    for res in st.session_state["results"]:
        name = res["name"]
        text = res["text"]
        meta = res["meta"]

        st.subheader(f"üìÑ {name}")
        with st.expander("PodglƒÖd", expanded=True):
            st.text_area(f"prev_{safe_filename(name)}", text, height=240, key=f"preview_{safe_filename(name)}")

        if isinstance(meta, dict) and meta.get("has_speakers"):
            st.markdown("##### üé§ Przypisz imiona m√≥wcom")
            new_text = speaker_mapper_form(name, res.get("original_text", text), text)
            if new_text is not None:
                res["text"] = new_text
                st.rerun()

    st.download_button(
        "‚¨áÔ∏è Pobierz wszystko w TXT",
        ("\n".join([f"\n{'='*60}\nPLIK: {r['name']}\n{'='*60}\n{r['text']}" for r in st.session_state['results']])).encode("utf-8"),
        file_name=f"converted_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
        mime="text/plain",
        key="dl_all_txt",
        disabled=st.session_state.get("converting", False)
    )

    c1, c2, c3, c4 = st.columns(4)

    with c1:
        if st.button("üíæ Zapisz po≈ÇƒÖczony TXT na dysk", key="btn_save_combined_txt",
                     disabled=st.session_state.get("converting", False)):
            out_dir = st.session_state.get("run_dir") or create_run_dir("outputs")
            st.session_state["run_dir"] = out_dir
            combined_path = os.path.join(out_dir, f"converted_{datetime.now().strftime('%Y%m%d_%H%M')}.txt")
            save_text(combined_path, st.session_state["combined_text"])
            st.success(f"Zapisano: {combined_path}")

    with c2:
        if st.button("üíæ Zapisz wszystkie SRT (audio)", key="btn_save_all_srt",
                     disabled=st.session_state.get("converting", False)):
            out_dir = st.session_state.get("run_dir") or create_run_dir("outputs")
            st.session_state["run_dir"] = out_dir
            saved = 0
            for it in st.session_state["results"]:
                meta = it.get("meta") or {}
                if meta.get("type") == "audio" and meta.get("segments"):
                    fname_base = os.path.splitext(safe_filename(it["name"]))[0]
                    srt_path = os.path.join(out_dir, f"{fname_base}.srt")
                    save_text(srt_path, segments_to_srt(meta["segments"]))
                    saved += 1
            st.success(f"Zapisano SRT dla {saved} plik√≥w audio w: {out_dir}")

    with c3:
        if st.button("üß† Generuj podsumowania (AI)", key="btn_make_summaries",
                     disabled=st.session_state.get("converting", False)):
            st.session_state["audio_summaries"] = []
            st.session_state["doc_summaries"] = []

            # AUDIO
            if summarize_audio_enabled and st.session_state["audio_items"]:
                for (aname, atext, ameta) in st.session_state["audio_items"]:
                    diarized = bool(ameta.get("has_speakers"))
                    with st.spinner(f"Tworzƒô podsumowanie AUDIO dla {aname}..."):
                        summary_json = summarize_meeting_transcript(
                            transcript=atext,
                            model=summarize_model if summarize_model_candidates else "llama3:latest",
                            max_chars=chunk_chars,
                            diarized=diarized
                        )
                        summary_md = build_meeting_summary_markdown(summary_json)
                        st.session_state["audio_summaries"].append(
                            {"name": aname, "md": summary_md, "json": summary_json}
                        )
            elif summarize_audio_enabled:
                st.info("Brak plik√≥w audio do podsumowania.")

            # DOKUMENTY TEKSTOWE / OBRAZY
            if summarize_text_enabled or summarize_images_enabled:
                any_doc = False
                for res in st.session_state["results"]:
                    meta = res.get("meta") or {}
                    t = meta.get("type")
                    is_text_doc = t in ("pdf", "docx", "pptx", "txt", "email")
                    is_image_doc = t == "image"
                    if ((is_text_doc and summarize_text_enabled) or
                        (is_image_doc and summarize_images_enabled)):
                        any_doc = True
                        name = res["name"]
                        text = res["text"]
                        with st.spinner(f"Tworzƒô podsumowanie DOKUMENTU dla {name}..."):
                            summary_json = summarize_document_text(
                                text=text,
                                model=summarize_model if summarize_model_candidates else "llama3:latest",
                                max_chars=chunk_chars,
                            )
                            summary_md = build_meeting_summary_markdown(summary_json)
                            st.session_state["doc_summaries"].append(
                                {"name": name, "md": summary_md, "json": summary_json}
                            )
                if not any_doc:
                    st.info("Brak dokument√≥w do podsumowania.")

            st.success("Gotowe podsumowania ‚Äì poni≈ºej do pobrania/zapisania.")

    with c4:
        if has_anythingllm:
            if st.button("üì§ Wy≈õlij do AnythingLLM", key="btn_send_anythingllm",
                         disabled=st.session_state.get("converting", False)):
                success, msg = send_to_anythingllm(st.session_state["combined_text"], "converted_docs.txt")
                if success:
                    st.success(msg)
                else:
                    st.error(msg)
        else:
            st.caption("AnythingLLM wy≈ÇƒÖczone lub brak config/tryb offline")

    # === WY≈öWIETLANIE PODSUMOWA≈É AUDIO ===
    if st.session_state.get("audio_summaries"):
        st.markdown("---")
        st.subheader("üß† Podsumowania rozm√≥w audio")

        for s in st.session_state["audio_summaries"]:
            aname = s["name"]
            summary_md = s["md"]
            summary_json = s["json"]

            st.markdown(f"### üéß {aname}")
            st.markdown(summary_md)

            col_a, col_b, col_c = st.columns(3)
            with col_a:
                st.download_button(
                    "‚¨áÔ∏è Pobierz MD",
                    summary_md.encode("utf-8"),
                    file_name=f"{safe_filename(aname).replace('.mp3', '').replace('.wav', '')}_summary.md",
                    mime="text/markdown",
                    key=f"dl_md_{safe_filename(aname)}",
                    disabled=st.session_state.get("converting", False)
                )
            with col_b:
                st.download_button(
                    "‚¨áÔ∏è Pobierz JSON",
                    json.dumps(summary_json, ensure_ascii=False, indent=2).encode("utf-8"),
                    file_name=f"{safe_filename(aname).replace('.mp3', '').replace('.wav', '')}_summary.json",
                    mime="application/json",
                    key=f"dl_json_{safe_filename(aname)}",
                    disabled=st.session_state.get("converting", False)
                )
            with col_c:
                if st.button("üíæ Zapisz (MD+JSON) na dysk", key=f"btn_save_sum_{safe_filename(aname)}",
                             disabled=st.session_state.get("converting", False)):
                    out_dir = st.session_state.get("run_dir") or create_run_dir("outputs")
                    st.session_state["run_dir"] = out_dir
                    base = safe_filename(aname).replace('.mp3', '').replace('.wav', '')

                    md_path = os.path.join(out_dir, f"{base}_summary.md")
                    json_path = os.path.join(out_dir, f"{base}_summary.json")

                    save_text(md_path, summary_md)
                    save_text(json_path, json.dumps(summary_json, ensure_ascii=False, indent=2))

                    st.success(f"Zapisano podsumowanie do: {out_dir}")

    # === WY≈öWIETLANIE PODSUMOWA≈É DOKUMENT√ìW ===
    if st.session_state.get("doc_summaries"):
        st.markdown("---")
        st.subheader("üß† Podsumowania dokument√≥w tekstowych / obraz√≥w")

        for s in st.session_state["doc_summaries"]:
            name = s["name"]
            summary_md = s["md"]
            summary_json = s["json"]

            st.markdown(f"### üìÑ {name}")
            st.markdown(summary_md)

            col_a, col_b, col_c = st.columns(3)
            with col_a:
                st.download_button(
                    "‚¨áÔ∏è Pobierz MD",
                    summary_md.encode("utf-8"),
                    file_name=f"{safe_filename(name)}_summary.md",
                    mime="text/markdown",
                    key=f"dl_md_doc_{safe_filename(name)}",
                    disabled=st.session_state.get("converting", False)
                )
            with col_b:
                st.download_button(
                    "‚¨áÔ∏è Pobierz JSON",
                    json.dumps(summary_json, ensure_ascii=False, indent=2).encode("utf-8"),
                    file_name=f"{safe_filename(name)}_summary.json",
                    mime="application/json",
                    key=f"dl_json_doc_{safe_filename(name)}",
                    disabled=st.session_state.get("converting", False)
                )
            with col_c:
                if st.button("üíæ Zapisz (MD+JSON) na dysk", key=f"btn_save_sum_doc_{safe_filename(name)}",
                             disabled=st.session_state.get("converting", False)):
                    out_dir = st.session_state.get("run_dir") or create_run_dir("outputs")
                    st.session_state["run_dir"] = out_dir
                    base = safe_filename(name)

                    md_path = os.path.join(out_dir, f"{base}_summary.md")
                    json_path = os.path.join(out_dir, f"{base}_summary.json")

                    save_text(md_path, summary_md)
                    save_text(json_path, json.dumps(summary_json, ensure_ascii=False, indent=2))

                    st.success(f"Zapisano podsumowanie do: {out_dir}")

# Reset sesji (nie rusza plik√≥w na dysku)
st.markdown("---")
if st.button("‚ôªÔ∏è Reset sesji (wyczy≈õƒá wyniki)", type="secondary", key="btn_reset_session",
             disabled=st.session_state.get("converting", False)):
    for k in ["results", "combined_text", "audio_items", "audio_summaries", "doc_summaries",
              "run_dir", "project_brain", "project_tasks"]:
        if isinstance(st.session_state.get(k), list):
            st.session_state[k] = []
        else:
            st.session_state[k] = None
    st.session_state["stats"] = {'processed': 0, 'errors': 0, 'pages': 0}
    st.session_state["converted"] = False
    st.info("Wyczyszczono wyniki z pamiƒôci sesji.")
