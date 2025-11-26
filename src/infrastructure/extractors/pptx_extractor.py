"""
PPTX Extractor - Extract text from PowerPoint presentations.

Strategy:
- Extract text from each slide
- Include notes if available
- Optionally use Vision for images
"""

import logging
from typing import BinaryIO
import base64

from pptx import Presentation

from domain.models.document import (
    ExtractionResult,
    Page,
    ExtractionMetadata,
    DocumentType
)
from domain.models.config import ExtractionConfig
from domain.interfaces.llm_client import VisionLLMClient
from domain.exceptions import ExtractionError

logger = logging.getLogger(__name__)


class PPTXExtractor:
    """
    PPTX presentation extractor with optional Vision support.
    """

    def __init__(self, vision_client: VisionLLMClient | None = None):
        """
        Initialize PPTX extractor.

        Args:
            vision_client: Optional vision model for image descriptions
        """
        self.vision = vision_client
        logger.info(f"PPTXExtractor initialized (vision={'enabled' if vision_client else 'disabled'})")

    @property
    def name(self) -> str:
        return "PPTX Extractor"

    @property
    def supported_extensions(self) -> tuple[str, ...]:
        return ('.pptx', '.ppt', '.PPTX', '.PPT')

    def can_handle(self, file_name: str) -> bool:
        """Check if file is a PPTX"""
        lower_name = file_name.lower()
        return lower_name.endswith('.pptx') or lower_name.endswith('.ppt')

    def extract(
        self,
        file: BinaryIO,
        file_name: str,
        config: ExtractionConfig
    ) -> ExtractionResult:
        """
        Extract text from PPTX.

        Args:
            file: PPTX file object
            file_name: Original file name
            config: Extraction configuration

        Returns:
            ExtractionResult with slides as pages

        Raises:
            ExtractionError: If extraction fails
        """
        import time
        start_time = time.time()

        try:
            file.seek(0)
            logger.info(f"Extracting PPTX: {file_name}")

            prs = Presentation(file)
            pages = []

            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== Slide {i} ==="]

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)

                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_content.append(f"\nNotes: {notes}")

                # Optional: Vision for images
                if self.vision and config.use_vision:
                    for shape in slide.shapes:
                        if getattr(shape, "shape_type", None) == 13:  # PICTURE
                            try:
                                img_stream = shape.image.blob
                                img_b64 = base64.b64encode(img_stream).decode()
                                description = self.vision.analyze_image(
                                    img_stream,
                                    config.vision_prompt,
                                    model=config.vision_model
                                )
                                slide_content.append(f"\n[Image]: {description}")
                            except Exception as e:
                                logger.warning(f"Vision failed for slide {i} image: {e}")

                slide_text = "\n".join(slide_content)
                pages.append(Page(number=i, text=slide_text))

            processing_time = time.time() - start_time

            metadata = ExtractionMetadata(
                document_type=DocumentType.PPTX,
                pages_count=len(pages),
                extraction_method="python-pptx" + (" + vision" if self.vision and config.use_vision else ""),
                processing_time_seconds=processing_time,
                file_size_bytes=self._get_file_size(file),
                vision_model=config.vision_model if (self.vision and config.use_vision) else None
            )

            logger.info(
                f"PPTX extraction complete: {file_name} | "
                f"{len(pages)} slides | {processing_time:.2f}s"
            )

            return ExtractionResult(
                file_name=file_name,
                pages=pages,
                metadata=metadata
            )

        except Exception as e:
            logger.exception(f"PPTX extraction failed: {file_name}")
            raise ExtractionError(
                f"Failed to extract PPTX: {e}",
                file_name=file_name
            ) from e

    @staticmethod
    def _get_file_size(file: BinaryIO) -> int:
        """Get file size in bytes"""
        current_pos = file.tell()
        file.seek(0, 2)
        size = file.tell()
        file.seek(current_pos)
        return size
